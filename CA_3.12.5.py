#!/usr/bin/env python3.12.6
# -*- coding: utf-8 -*-

#CORals Analytics v.3.12.5
#This script is created by Christian Otto Ruge and CORals.
#It is licenced under GNU GPL v.3.
#https://www.corals.no


import os
from turtle import fd
import pandas as pd
import numpy as np
import PySimpleGUI as sg 
import xlrd
import xlsxwriter
from pyprocessmacro import Process
from io import open
import openpyxl
import pingouin as pg
from pingouin import mediation_analysis
import seaborn as sns
import matplotlib.pyplot as plt
from sklearn import linear_model
from scipy import stats
from scipy.stats import zscore
from scipy.stats import skew
from scipy.stats import kurtosis
import statsmodels.api as sm
from statsmodels.stats.outliers_influence import variance_inflation_factor
import warnings
from factor_analyzer import FactorAnalyzer
from factor_analyzer.factor_analyzer import calculate_bartlett_sphericity
from factor_analyzer.factor_analyzer import calculate_kmo
from factor_analyzer import (ConfirmatoryFactorAnalyzer, ModelSpecificationParser)
import pyreadstat as py
from string import printable
import datetime
from datetime import datetime
from datetime import date
from pptx import Presentation
from pptx.util import Inches
import sys
import tracemalloc
#import df2img
import dataframe_image as dfi
import json
from stepmix.stepmix import StepMix
from stepmix.utils import get_mixed_descriptor

tracemalloc.start()

warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=FutureWarning)

warnings.simplefilter("ignore", category=FutureWarning)

# From ChatGPT 10.10.23
# GitHUB: https://github.com/keita43a/Postal_code_region_number_Norway

# Reading file within pyinstaller exe-file

if getattr(sys, 'frozen', False):
    application_path = os.path.dirname(sys.executable)
elif __file__:
    application_path = os.path.dirname(__file__)
    #application_path = ''

def my_path(path_name):
    """Return the appropriate path for data files based on execution context"""
    if getattr(sys, 'frozen', False):
        # running in a bundle
        return(os.path.join(sys._MEIPASS, path_name))
        # NB: _MEI39442
    else:
        # running live
        return path_name   

CA_logo_file = 'Corals_Analytics_std.png'
logo_file_pyinstaller = my_path(CA_logo_file)
CA_logo_full_path = os.path.join(application_path, logo_file_pyinstaller)

def df_save_to_image(df, filename):
    df = df.copy()
    df = df.style.set_table_styles([dict(selector='th', props=[('text-align', 'center'),('background-color', '#40466e'),('color', 'white')])])
    df.set_properties(**{'text-align': 'center'}).hide(axis='index')
    pd.set_option('colheader_justify', 'center')
    print(df)
    #html = df.to_html()
    #imgkit.from_string(html, filename)
    dfi.export(df, filename)

#Popup for exceptions:
def popup_break(event, message):
    while True:
        sg.Popup(message)
        event=='False'
        break

def CheckInt(s):
    try: 
        int(s)
        return True
    except ValueError:
        return False

#def EditDatasetSyntax (file):
#    from file import syntax_df
#    for row in syntax_df:
#        data_in[row[0]]=data_in[row[1]]

today = date.today()
date_time_str = today.strftime("%Y-%m-%d_%H-%M") 

def CreatePresentation(ev, val):

    global logo
    logo = CA_logo_full_path
    global filename_presentation
    global presentation_output_folder

    if (ev=='Continue') and val['presentation_old'] == True and val['presentation_file'] == '*Choose existing presentation':
        popup_break(ev, 'Choose existing presentation to add slides to')

    elif ev=='Continue' and val['presentation_new'] == True:
        print('Creating new presentation')
        global prs
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = 'CORals Analytics'
        subtitle.text = 'Results'         
        left_logo = Inches(8.0)
        top_logo = Inches(6.5)
        height_logo = Inches(0.8)
        pic = slide.shapes.add_picture(logo, left_logo, top_logo, height=height_logo)
        filename_presentation = 'CORals_Presentation_' + date_time_str + '.pptx'
        #filename_presentation = os.path.join(output_folder, file_presentation)
        presentation_output_folder = output_folder  

    elif ev=='Continue' and val['presentation_old'] == True:
        print('Detecting existing presentation')
        old_filename = val['presentation_file']
        prs = Presentation(old_filename)
        presentation_output_folder = output_folder
        #presentation_output_folder = os.path.abspath(old_filename)
        filename_presentation = os.path.basename(old_filename)
        print('Existing presentation detected')

    print('Filename presentation: ' + filename_presentation)

def AddToPresentation(image_output_folder, image_filename):
    print('Adding slide to presentation')
    left_logo = Inches(8.0)
    top_logo = Inches(6.5)
    height_logo = Inches(0.8)
    blank_slide_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_slide_layout)
    image_fullpath = os.path.join(image_output_folder, image_filename)
    pic = new_slide.shapes.add_picture(logo, left_logo, top_logo, height=height_logo)
    im = plt.imread(image_fullpath)
    h, w, c = im.shape
    print('Ratio h/b: ' + str(h) + '/' + str(w))
    
    if w/h > 1.33: #Landscape 
        n = w/h
        left = Inches(1.0)
        width = Inches(8)
        top = Inches(0.5 + (3-(3.5/n))) #Half of page hight
        pic = new_slide.shapes.add_picture(image_fullpath, left, top, width=width)

    elif w/h < 1: #Vertical NB: Incomplete!!
        n = h/w
        left = Inches(1 + (4 -(3/n))) #Half of page width
        height = Inches(6.0)
        top = Inches(0.5)
        pic = new_slide.shapes.add_picture(image_fullpath, left, top, height = height)
        
    else:
        n = w/h #Landscape or square NB: Incomplete!!
        left = Inches(1 + (2 -(n/4))) #Half of page width
        height = Inches(6.0)
        top = Inches(0.5)
        pic = new_slide.shapes.add_picture(image_fullpath, left, top, height = height)

    pic = new_slide.shapes.add_picture(logo, left_logo, top_logo, height=height_logo)

def SavePresentation(val):

    if (val['presentation_old'] == True or val['presentation_new'] == True):
        prs.save(os.path.join(presentation_output_folder, filename_presentation))
        print('Presentation saved')

sg.theme('LightGrey1')

layoutOriginal = [
    [sg.Text('CORals analytics', size=(25,1), justification='left', font=("Arial", 20, "bold"))],
    [sg.Text('TOOLS FOR SCIENTIFIC RESEARCH', font=('bold'))],
    [sg.Text('')],  
    [sg.Text('Choose a CORals-tool', key='Choose', font=('bold'))],     
    [sg.Frame(layout=[      
    [sg.Radio('Distribution analysis', "RADIO1", key="Distribution", default=False, size=(18,1)), sg.Radio('Correlation analysis', "RADIO1", key="Correlation",  default=False, size=(18,1)), sg.Radio('Regression analysis', "RADIO1", key="Regression", default=False, size=(23,1))],         
    [sg.Radio('Mediation analysis', "RADIO1", key="Mediation", default=False, size=(18,1)), sg.Radio('Moderation analysis', "RADIO1", key="Moderation", default=False, size=(18,1)), sg.Radio('Factor analysis', "RADIO1", key="Factor", default=False, size=(23,1))],
    [sg.Radio('Cronbachs Alpha', "RADIO1", key="Cronbachs", default=False, size=(18,1)), sg.Radio('CSV Rescue', "RADIO1", key="CSV", default=False, size=(18,1)), sg.Radio('Latent Class Analysis', "RADIO1", key="lca", default=False, size=(23,1))],
    [sg.Radio('Edit dataset', "RADIO1", key="Compute", default=False, size=(18,1)), sg.Radio('SAV-converter', "RADIO1", key="sav", default=False, size=(18,1)),sg.Radio('Make presentation', "RADIO1", key="Presentation", default=False, size=(23,1))]], title='Options', title_color='red', relief=sg.RELIEF_SUNKEN, tooltip='Use these to set flags')],    
    [sg.Text('Choose dataset (except for "CSV Rescue-"/"SAV-converter")', font=('bold'))],
    [sg.Text('NB: Dataset must be in ".xlsx" or ".csv". All variables to be analysed should be numeric.)')],
    [sg.In('', key='analytics_dataset', size=(70,1)), sg.FileBrowse()],
    [sg.Text('(Optional) To reduce dataset by variable and value, enter variable and value: ')],
    [sg.Text('Variable: '), sg.In('', key="filter_var", size=(10,1)), sg.Text(' and value: '), sg.In('', key='filter_val', size=(5,1))],
    [sg.Text('(Optional) Choose syntax file (.py) to edit dataset temporarily:')],
    [sg.In('', key='syntax_file', size=(70,1)), sg.FileBrowse()],
    [sg.Button('Continue'), sg.Button('Close'), sg.Button('Documentation')],
    [sg.Text('')],
    [sg.Text('', size=(19,2)), sg.Text('www.corals.no', font=("Ariel", 20, 'bold'), text_color=('lightsteelblue2'))]]
    #Logo excluded for pyinstaller to work properly:
    #[sg.Column([[sg.Image('logo.png', key='icon', size=(340, 150))]], justification='center')]
    


layoutOriginalFrame = [
    [sg.Frame('Functions', layoutOriginal, size= (570, 615))]]

columnRight = [
    [sg.Text('A brief introduction:', font=('bold'), size = (45, 1), justification='left', background_color='lightsteelblue2')],  
    [sg.Text('"Distribution analysis" outputs descriptive statistics (distribution values). Charts, bars, pairwise plots and values may be exported (optional).', size=(45,3), background_color='lightsteelblue2')],
    [sg.Text('"Correlation analysis" outputs correlation values and p-values. Correlation matrix and heatmap may be exported (optional).', size=(45,2), background_color='lightsteelblue2')],
    [sg.Text('"Regression analysis" outputs simple, or multiple, regression values and VIF values. Plots and tables may be exported (optional).', size=(52,2), background_color='lightsteelblue2')],
    [sg.Text('"Mediation analysis" outputs mediation values such as "direct effect" and "indirect effect". Tables may be exported (optional).', size=(45,2), background_color='lightsteelblue2')],
    [sg.Text('"Moderation analysis" is doing moderation analyses according to Dr. F. Hayes (model 1) and outputs moderation values. Line charts and tables may be exported (optional).', size=(45,3), background_color='lightsteelblue2')],
    [sg.Text('"Factor analysis" performs exploratory analysis, confirmatory analyses and other factor analyses. Values may be exported (optional).', size=(52,2), background_color='lightsteelblue2')],  
    [sg.Text('"Edit dataset" creates new variables in the dataset (scale-variables or dictonomous variables), and deletes unwanted variables.', size=(45,3), background_color='lightsteelblue2')],  
    [sg.Text('"Cronbachs Alpha" is a validation tool for factors.', size = (45, 2), justification='left', background_color='lightsteelblue2')],
    [sg.Text('"CSV rescue" converts csv-files to xlsx-files, resolving encoding problems (ascii to latin) during the process.', size = (45, 2), justification='left', background_color='lightsteelblue2')],
    [sg.Text('"SAV-converter" converts sav-files (SPSS) to both csv- and xlsx-files', size = (45, 1), justification='left', background_color='lightsteelblue2')],
    [sg.Text('"Make presentation" creates a Powerpoint presentation with the results.', size = (45, 1), justification='left', background_color='lightsteelblue2')],
    [sg.Text('"Documentation" gives information about which Python tools that are used for each tool.', size=(45,2), background_color='lightsteelblue2')]]

layoutRight = [
    [sg.Frame('Description', columnRight, size=(430,615), background_color='lightsteelblue2')]]  
    
layoutBottom =  [
    [sg.Text('')],
    [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

layoutComplete = [
    [sg.Column(layoutOriginalFrame, element_justification='left', size=(580,620)), sg.Column(layoutRight, element_justification='right', size=(505,620))], 
    [sg.Column(layoutBottom, element_justification='left')]]

try:    
    winOriginal = sg.Window('CORals Analytics v. 3.12.5', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutComplete)    

    winCsv_active=False
    winCorrelation_active=False
    winRegression_active=False
    winMediation_active=False
    winModeration_active=False
    winDistribution_active=False
    winFactor_active=False
    winCompute_active=False
    winFactor_active=False
    winDocumentation_active=False
    winSav_active=False
    winCronbachs_active=False
    winPresentation_active=False
    winLca_active=False

    while True:
        evOriginal,valOriginal=winOriginal.Read(timeout=100)
        if evOriginal is None or evOriginal=='Close':
            winOriginal.Close()
            del winOriginal
            break

        #Reveal datatype and create dataframe
        data_in=valOriginal['analytics_dataset']
        data_in_str=str(data_in)
        chosen_filename=os.path.basename(data_in)
        #syntax_file=valOriginal['syntax_file']

        if (evOriginal)=='Continue':
            if data_in.endswith('.csv'):
                data_df=pd.read_csv(data_in, sep=None, encoding='iso-8859-1', engine='python')
                if not (valOriginal['filter_var']=='' and valOriginal['filter_val']==''):
                    filter_var=valOriginal['filter_var']
                    filter_val=valOriginal['filter_val']
                    filter_val=int(filter_val)
                    if filter_var in data_df:
                        data_df=data_df.loc[data_df[filter_var]==filter_val]
            
            elif data_in.endswith('.xlsx'):
                data_df=pd.read_excel(data_in, engine='openpyxl')
                if not (valOriginal['filter_var']=='' and valOriginal['filter_val']==''):
                    filter_var=valOriginal['filter_var']
                    filter_val=valOriginal['filter_val']
                    filter_val=int(filter_val)
                    if filter_var in data_df:
                        data_df=data_df.loc[data_df[filter_var]==filter_val]

        if (evOriginal)=='Continue' and (not ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx')))) and not ((valOriginal['Correlation']==False) and(valOriginal['lca']==False) and (valOriginal['Regression']==False) and (valOriginal['Mediation']==False) and (valOriginal['Moderation']==False) and (valOriginal['Distribution']==False) and (valOriginal['Compute']==False) and (valOriginal['Factor']==False) and (valOriginal['Cronbachs']==False)):
            popup_break(evOriginal, 'Please enter a valid dataset.')

        if (evOriginal)=='Continue' and (((data_in.endswith('.csv')) and valOriginal['Compute']==True)) and not ((valOriginal['Correlation']==False) and (valOriginal['Regression']==False) and (valOriginal['Mediation']==False) and (valOriginal['Moderation']==False) and (valOriginal['Distribution']==False) and (valOriginal['Compute']==False) and (valOriginal['Factor']==False) and (valOriginal['Cronbachs']==False)):
            popup_break(evOriginal, '"Edit variable" is only available for .xlsx-files.')

        if (not winCsv_active) and (valOriginal['CSV']==True) and (evOriginal=='Continue'):
            winOriginal.Hide()
            winCsv_active=True         

            layoutCsv = [
                [sg.Text('')], 
                [sg.Text('CSV rescue:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('')], 
                [sg.Text('Choose CSV-file:', font=('bold'))],
                [sg.In('', key='csv-file'), sg.FileBrowse()],
                [sg.Text('Select output folder:', size=(35, 1), font=('bold'))],      
                [sg.InputText('', key='xlsx_output'), sg.FolderBrowse()],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Button('Fix and convert'), sg.Button('Back')]]
            
            layoutOriginalFrame = [
                [sg.Frame('Function', layoutCsv, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('', background_color='lightsteelblue2')],
                [sg.Text('"CSV rescue" converts csv-files to xlsx-files, resolving encoding problems (ascii to latin) during the process.', background_color='lightsteelblue2', size = (50, 2), justification='left')], 
                [sg.Output(size=(60,2))]], background_color='lightsteelblue2', size=(500,600))]]
            
            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteCsv = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]
            
            winCsv=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteCsv)

            while True:
                evCsv, valCsv = winCsv.Read(timeout=100)
                
                if evCsv is None or evCsv == 'Back':
                    winCsv_active=False
                    winCsv.Close()
                    del winCsv
                    winOriginal.UnHide()
                    break
                
                if (evCsv=='Fix and convert') and (valCsv['csv-file']==''):
                    popup_break(evCsv, 'Choose csv-file')

                if (evCsv=='Fix and convert') and (valCsv['xlsx_output']==''):
                    popup_break(evCsv, 'Choose output folder')

                if not (valCsv['csv-file']=='') and not (valCsv['xlsx_output']=='') and (evCsv=='Fix and convert'):   
                    while True:
                        dataset = valCsv['csv-file']
                        filename=os.path.basename(dataset)
                        csv_file=str(filename)
                        folder=dataset.replace(csv_file,'')
                        output_folder=valCsv['xlsx_output']
                        df=pd.read_csv(dataset, sep=None, encoding='iso-8859-1', engine='python')
                        xlsx_filename=csv_file.replace(".csv", ".xlsx")

                        engine = 'xlsxwriter'

                        with pd.ExcelWriter(os.path.join(output_folder, xlsx_filename), engine=engine) as writer:
                            df.to_excel(writer, sheet_name="Edited", index = None, header=True)
                            
                        print('Done: Csv-file fixed and converted to xlsx-file')

                        break
            
        #Descriptive/Distribution
        if (not winDistribution_active) and (valOriginal['Distribution']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winDistribution=True


            layoutDistribution = [
                [sg.Text('Distribution analysis:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('Enter variables to be analyzed:', font=('bold'))],
                [sg.Text('Divide included variables by commas.')],
                [sg.Text('(Variable names are the column headers (row 1) in the dataset)')],
                [sg.InputText('', key='distribution_variables', size=(70,1))],
                [sg.Text('Export output: ', font=('bold')), sg.Radio('No: ', 'RADIO2', key='distribution_output_no', default=True, size=(20,1)), sg.Radio('Yes: ', 'RADIO2', key='distribution_output_yes', default=False, size=(20,1))],
                [sg.Text('Select output formats:', font=('bold'))],
                [sg.Checkbox('Line plots', key='lines', default=False, size=(12,1)),  sg.Checkbox('Bar plots', key='bars', default=False, size=(12,1)), sg.Checkbox('Pairwise plots', key='plots', default=False, size=(12,1)), sg.Checkbox('Descriptive values', key='values', default=False, size=(15,1))],  
                [sg.Radio('No presentation', 'RADIO3', key='no_presentation', default=True, size=(12,1)), sg.Radio('New presentation', 'RADIO3', key='presentation_new', default=False, size=(12,1)), sg.Radio('Add to presentation*', 'RADIO3', key='presentation_old', default=False, size=(15,1))], 
                [sg.InputText('*Choose existing presentation', key='presentation_file', size=(60,1)), sg.FileBrowse()],
                [sg.Text('Optional: Enter hue-variables (for pairwise plots only)')],
                [sg.Text('Hue variable must be one of the above selected variables.')],
                [sg.InputText('', key='hue', size=(15, 1))],
                [sg.Text('Select output folder:', font=('bold'), size=(60, 1))],      
                [sg.InputText('', key='distribution_output', size=(60,1)), sg.FolderBrowse()],
                [sg.Button('Continue'), sg.Button('Back')]]
            
            layoutLeft=[    
                [sg.Frame('Function', layoutDistribution, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,600), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteDistribution = [[sg.Column(layoutLeft, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]
            
            winDistribution=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteDistribution)
            
            while True:
                evDistribution, valDistribution = winDistribution.Read(timeout=100)
  
                if evDistribution=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ', '
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evDistribution==''
                
                if evDistribution in (None, 'Back') or evDistribution =='Step back and change dataset':
                    winDistribution_active=False
                    winDistribution.Close()
                    del winDistribution
                    winOriginal.UnHide()
                    break       
                    
                elif (evDistribution=='Continue' and valDistribution['distribution_variables']==''):
                    popup_break(evDistribution, 'Choose variables')

                elif (evDistribution=='Continue') and (not valDistribution['hue']=='') and ((valDistribution['lines']==True) or (valDistribution['values']==True) or (valDistribution['bars']==True)) and (valDistribution['plots']==False):
                    popup_break(evDistribution, 'Enter hue value only for "pairwise plots"')

                elif (evDistribution=='Continue') and valDistribution['distribution_output_yes']==True and (valDistribution['lines']==False) and (valDistribution['bars']==False) and (valDistribution['plots']==False) and (valDistribution['values']==False):
                    popup_break(evDistribution, 'Choose output format')

                elif (evDistribution=='Continue') and (valDistribution['distribution_output_yes']==True) and ((valDistribution['presentation_old'] == True) or (valDistribution['presentation_new'] == True)) and (valDistribution['lines']==False) and (valDistribution['bars']==False) and (valDistribution['plots']==False):
                    popup_break(evDistribution, 'Please select graphics for the presentation')

                elif (evDistribution=='Continue') and (valDistribution['distribution_output_yes']==True) and (valDistribution['presentation_old'] == True) and not valDistribution['presentation_file'].endswith('.pptx'):
                    popup_break(evDistribution, 'Existing presentation has to be a Microsoft Powerpont file (.pptx)')

                elif (evDistribution=='Continue') and (',' in valDistribution['hue']):
                    popup_break(evDistribution, 'select only one hue variable')
                
                elif (evDistribution=='Continue') and ((not valDistribution['hue']=='') and (valDistribution['plots'] == False)):
                    popup_break(evDistribution, 'Hue variable is only available for pairwise plots.')                                                                           

                if (evDistribution=='Continue') and not (valDistribution['hue']==''):
                    color=str(valDistribution['hue'])                   

                elif (evDistribution=='Continue') and (valDistribution['hue']==''):
                    color=None
                    
                #Valdidate variables supplied
                if (evDistribution=='Continue'):

                    data=data_df

                    validvariables=list(data.columns)
                    var=valDistribution['distribution_variables']
                    var=var.replace(" ", "")
                    varnames=str(var)
                    varnames=varnames.replace('[','')
                    varnames=varnames.replace(']','')
                    varsplit=var.split(',')
                    list_var=list(varsplit)
                    response='yes'
               
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'

                    if (evDistribution=='Continue') and (not valDistribution['hue']=='') and (not valDistribution['hue'] in str(list_var)):   
                        popup_break(evDistribution, 'Hue variable must be one of the above entered variables.')
                    
                    elif (evDistribution=='Continue') and (response =='no'):
                        popup_break(evDistribution, 'Variable(s) not in dataset')              
                    
                    elif (evDistribution=='Continue') and (valDistribution['distribution_output']=='' and (valDistribution['distribution_output_yes']==True)):
                        popup_break(evDistribution, 'Choose output folder')   
                    

                    elif (evDistribution=='Continue') and (response=='yes'):                                                
                           
                        while True:

                            df=data[list_var]
                            values=df.describe().round(2)
                            print('')
                            print('Values:')
                            print('')
                            print(values)
                            print('')
                            
                            for n in list_var:
                                var=n
                                count = df[n].value_counts().reset_index()
                                norm_df = df[n].value_counts(normalize=True).reset_index()
                                norm_df = norm_df.round(3)
                                
                                total_df=count.join(norm_df['proportion'])
                                total_df.columns=['Variable', 'Count', 'Percent']

                                for i, row in total_df.iterrows():
                                    total_df.at[i,'Percent'] = str(round((total_df.iloc[i,2] * 100),1)) + "%"
                                
                                
                                print(f'Distribution of: {n}')
                                print(total_df.to_string(index=False))
                                print('')

                                skew_pandas=df[n].skew()
                                kurt_pandas=df[n].kurtosis()
                                print(f'Skew: {skew_pandas}')
                                print(f'Kurtosis: {kurt_pandas}')
                                print('')
                                
                                output_folder = valDistribution['distribution_output']

                                if valDistribution['values']==True and valDistribution['distribution_output_yes']==True:
                                    engine = 'xlsxwriter'
                                    new_filename=f'Distribution_{n}.xlsx'
                                    
                                    #values = df[n].tolist()
                                    #skew_value = skew(values, axis=0, bias=False)
                                    #kurtosis_value = kurtosis(values, axis=0, bias=False)
                                    empty_list=[]
                                    skew_kurt_df = pd.DataFrame(empty_list)
                                    skew_kurt_df[0] = ['Variable: ', 'Scew: ', 'Kurtosis: ']
                                    skew_kurt_df[1] = [n, skew_pandas, kurt_pandas]
                                                                        
                                    #count = count.rename(columns={count.columns[0]: 'Value'})
                                    with pd.ExcelWriter(os.path.join(output_folder, new_filename)) as writer:                     
                                        values.to_excel(writer, sheet_name="Distribution", startcol=0, startrow=0)
                                        total_df.to_excel(writer, sheet_name="Distribution", index=False, startcol=0, startrow=10)
                                        skew_kurt_df.to_excel(writer, sheet_name="Distribution", index=False, header=False, startcol=0, startrow=18)
                                    print(f'Variable: {n} - values saved to file')
                                    print('')

                            print('')
                            print('Distribution analysis completed')
                            print('')
                            
                            if valDistribution['distribution_output_yes']==True:
                                if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                    CreatePresentation(evDistribution, valDistribution)                


                                if valDistribution['lines']==True:

                                    # Single plots:
                                    for n in list_var:

                                        df[n].plot(kind="kde", title=n, ylabel="Distribution", xlim=(df[n].min(),df[n].max()))
                                        
                                        new_filename="Lineplot_" + n + "_distribution.png"
                                        
                                        print('New filename: ' + new_filename)
                                        plt.savefig(os.path.join(output_folder, new_filename), dpi=300, format='png', transparent=True)
                                        #plt.show()
                                        plt.close()
                                        print(f'Line plots for {n}saved')
                                        print('')

                                        if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                            AddToPresentation(output_folder, new_filename)

                                    ''' Requires same alternatives/scales
                                    # Combined plot:
                                    sns.lineplot(x='Alternatives', y='Score', hue='Variables', data=pd.melt(df, [list_var]))
                                    
                                    new_filename_combined="Lineplot_combined" + n + "_distribution.png"
                                    
                                    print('New filename: ' + new_filename)
                                    plt.savefig(os.path.join(output_folder, new_filename), dpi=300, format='png', transparent=True)
                                    #plt.show()
                                    plt.close()
                                    print(f'Combined line plot for {n}saved')
                                    print('')

                                    if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                        AddToPresentation(output_folder, new_filename_combined)
                                    '''

                                if valDistribution['bars']==True:
                                
                                    for n in list_var:

                                        df[n].plot(kind="hist", edgecolor="black", x="Values", y="Numbers", title=n, ylabel="Distribution")
                                        
                                        new_filename="Bars_" + n + "_distribution.png"
                                        
                                        print('New filename: ' + new_filename)
                                        plt.savefig(os.path.join(output_folder, new_filename), dpi=300, format='png', transparent=True)
                                        #plt.show()
                                        plt.close()
                                        print(f'Bar plot for {n}saved')
                                        print('')
                                        
                                        if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                            AddToPresentation(output_folder, new_filename)


                                if valDistribution['plots']==True:
                                    sns.pairplot(df, diag_kind="auto", hue=color, palette="husl")
                                    #plt.show()
                                    new_filename_png=f'Pairwise_plots_{varnames}.png'
                                    new_filename_pdf=f'Pairwise_plots_{varnames}.pdf'
                                    plt.savefig(os.path.join(output_folder, new_filename_png), dpi=300, format='png', transparent=True) 
                                    plt.savefig(os.path.join(output_folder, new_filename_pdf), dpi=300, format='pdf') 
                                    
                                    print(f'Pairwise plots for {varnames} saved')
                                    plt.close()

                                    if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                        filepath_png = os.path.join(output_folder, new_filename_png)
                                        AddToPresentation(output_folder, new_filename_png)

                                if valDistribution ['presentation_new'] == True or valDistribution['presentation_old'] == True:
                                    SavePresentation(valDistribution)

                            print('...')
                            break

        #Correlation
        
        if (not winCorrelation_active) and (valOriginal['Correlation']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winCorrelation_active=True


            layoutCorrelation = [
                [sg.Text('Bivarate correlation analysis:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('')],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('Enter variables to be analyzed:', font=('bold'))],
                [sg.Text('(Variable names are the column headers (row 1) in the dataset file.)')],
                [sg.Text('Divide variables by commas only (no white space).')],
                [sg.InputText('', key='correlation_variables')],
                [sg.Text('')],
                [sg.Text('Export output: ', font=('bold')), sg.Radio('No: ', 'RADIO1', key='corr_output_no', default=True, size=(10,1)), sg.Radio('Yes: ', 'RADIO1', key='corr_output_yes', default=False, size=(20,1))],
                [sg.Text('Select output formats:', font=('bold'))],
                [sg.Checkbox('Correlation matrix', key='matrix', default=False, size=(15,1)),  sg.Checkbox('Correlation heatmap', key='heatmap', default=False, size=(20,1))],  
                [sg.Radio('No presentation', 'RADIO3', key='no_presentation', default=True, size=(15,1)), sg.Radio('New presentation', 'RADIO3', key='presentation_new', default=False, size=(15,1)), sg.Radio('Add to presentation*', 'RADIO3', key='presentation_old', default=False, size=(15,1))], 
                [sg.InputText('*Choose existing presentation', key='presentation_file', size=(60,1)), sg.FileBrowse()],
                [sg.Text('Select output folder:', font=('bold'), size=(35, 1))],      
                [sg.InputText('', key='correlation_output'), sg.FolderBrowse()],
                [sg.Button('Continue'), sg.Button('Back')],
                [sg.Text('')],
                [sg.Text('')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutCorrelation, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('"Correlation analysis" produces a correlation matrix or a graphical visual heatmap. Heatmap is only available for export.', size=(50,2), background_color='lightsteelblue2')], 
                [sg.Text('"The heatmap displays correlation values only, while the matrix also displays significance values.', size=(50,2), background_color='lightsteelblue2')]], size=(500,100), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,500), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteCorrelation = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winCorrelation=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteCorrelation)

            while True:
                evCorrelation, valCorrelation = winCorrelation.Read(timeout=100)    
                today = date.today()
                now = today.strftime("%Y-%m-%d")


                if evCorrelation=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evCorrelation==''

                if evCorrelation is None or evCorrelation=='Back' or evCorrelation=='Step back and change dataset':
                    winCorrelation_active=False
                    winCorrelation.Close()
                    del winCorrelation
                    winOriginal.UnHide()
                    break       

                    
                elif (evCorrelation=='Continue' and valCorrelation['correlation_variables']==''):
                    popup_break(evCorrelation, 'Choose variables')


                #Valdidate variables supplied
                elif (evCorrelation=='Continue') and not (valCorrelation['correlation_variables']==''):

                    var=valCorrelation['correlation_variables']
                    var = ''.join(char for char in var if char in printable)
                    var=var.replace(" ", "")
                    varsplit=var.split(',')

                    data=data_df[varsplit]
                    #remove rows with empty cells
                    data = data.dropna(axis=0)
                    #Remove rows with non-integar data
                    #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                    #Remove rows with non-numeric data
                    data = data.apply(pd.to_numeric, errors='coerce')

                    validvariables=list(data.columns)
                    
                    var=valCorrelation['correlation_variables']
                    varsplit=var.split(',')
                    data=data_df[varsplit]
                    list_var=list(varsplit)

                    var_tofilename=var.replace(',', '_')
                    var_tofilename_len=len(var_tofilename)
                    if var_tofilename_len > 50:
                        var_tofilename = var_tofilename[0:50] + '...'

                    now=datetime.today().strftime('%Y-%m-%d_%H-%M-%S')

                    response='yes'
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'

                    if (evCorrelation=='Continue') and (response =='no'):
                        popup_break(evCorrelation, 'Variable(s) not in dataset')

                    #End validation of variables
                    
                    elif (evCorrelation=='Continue' and valCorrelation['corr_output_yes']==True and (valCorrelation['heatmap']==False and valCorrelation['matrix']==False)):
                        popup_break(evCorrelation, 'Choose output format')  
                    
                    elif (evCorrelation=='Continue' and valCorrelation['corr_output_no']==True and (valCorrelation['heatmap']==True)):
                        popup_break(evCorrelation, 'Heatmap is only available for ouput.')  

                    elif (evCorrelation=='Continue' and valCorrelation['corr_output_yes']==True and valCorrelation['correlation_output']==''):
                        popup_break(evCorrelation, 'Choose output folder')   

                    elif (evCorrelation=='Continue') and (valCorrelation['corr_output_yes']==True) and ((valCorrelation['presentation_old'] == True) or (valCorrelation['presentation_new'] == True)) and (valCorrelation['matrix']==False) and (valCorrelation['heatmap']==False):
                        popup_break(evCorrelation, 'Please select graphics for the presentation')

                    elif (evCorrelation=='Continue') and (valCorrelation['corr_output_yes']==True) and (valCorrelation['presentation_old'] == True) and not valCorrelation['presentation_file'].endswith('.pptx'):
                        popup_break(evCorrelation, 'Existing presentation has to be a Microsoft Powerpont file (.pptx)')

                    elif (evCorrelation=='Continue') and not (valCorrelation['correlation_variables']=='') and (response=='yes'):

                        while True:
                            output_folder=valCorrelation['correlation_output']
                            output_folder=os.path.abspath(output_folder)
                            
                            chosendata=data
                            #Create and correlation matrix
                            plot=pd.DataFrame(chosendata.rcorr(stars=True))
                            p_info='*=p<0.05 **=p<0.01 ***=p<0.001'
                            description=pd.DataFrame([['CORRELATION MATRIX:'],
                            [''],
                            ['Upper triangle: p-values'],
                            [p_info],
                            [''],
                            ['Lower triangle: Pearson r-values:'],
                            ['']])
                            
                            print(plot)
                            print('')
                            p_info='*=p<0.05 **=p<0.01 ***=p<0.001'
                            print(p_info)
                            print('')
                            print('...')

                            
                            if valCorrelation['corr_output_yes']==True:

                                if valCorrelation['presentation_new'] == True or valCorrelation['presentation_old'] == True:
                                    CreatePresentation(evCorrelation, valCorrelation) 

                                if valCorrelation['heatmap']==True and valCorrelation['corr_output_yes']==True and valCorrelation['correlation_output'] !='':
                                    print('Creating heatmap')
                                    #Create and save heatmap as png-file
                                    corrs = chosendata.corr()
                                    #mask = np.zeros_like(corrs)
                                    #mask[np.triu_indices_from(mask)] = True
                                    plt.figure()                                    
                                   
                                    heatmap=sns.heatmap(corrs, cmap='Spectral_r', annot=False, square=True, vmin=-1, vmax=1)
                                    heatmap.set(xlabel="", ylabel="")
                                    heatmap.set_title('Correlation Heatmap', fontdict={'fontsize':15}, pad=12);
                                    heatmap.invert_yaxis()
                                    heatmap.set_xticklabels(heatmap.get_xticklabels(), rotation=90, horizontalalignment='center')
                                    
                                    for tick in heatmap.get_yticklabels():
                                        tick.set_rotation(0)


                                    #varnames=str(varsplit)
                                    #defname=varnames.replace('[','')
                                    #defname=varnames.replace(']','')
                                    #plt.title('CORRELATION-MATRIX (HEATMAP):')
                                    #b, t = plt.ylim() # discover the values for bottom and top (Thanks to SalMac86, GitHub, 25.10.19)
                                    #b += 0.5 # Adds 0.5 to the bottom
                                    #t -= 0.5 # Subtracts 0.5 from the top
                                    #plt.ylim(b, t) # updates the ylim(bottom, top) values
                                    #plt.tight_layout()
                                    #plt.show()
                                  
                              
                                    filename = 'Correlation-heatmap_' + var_tofilename + '_' + now +'.png'
                                    new_filename=str(os.path.join(output_folder, filename))
                                    plt.savefig(new_filename, dpi=300, bbox_inches='tight')
                                    #plt.close(fig)

                                    if valCorrelation['presentation_new'] == True or valCorrelation['presentation_old'] == True:
                                        AddToPresentation(output_folder, filename)

                                    fig=None
                                    print('')
                                    print('Heatmap saved')
                                    plt.close()

                                if valCorrelation['matrix']==True and valCorrelation['corr_output_yes']==True and valCorrelation['correlation_output'] !='':
                                    plot_df = plot
                                    list_variable_names=plot.columns.values.tolist() 
                                    plot_df.insert(loc=0, column='Variables', value=list_variable_names)
                                    engine = 'xlsxwriter'
                                    with pd.ExcelWriter(os.path.join(output_folder, 'Correlation-plot_') + '_' + var_tofilename + '_' + now + '.xlsx', engine=engine) as writer:
                                        description.to_excel(writer, sheet_name="pearson r and p-values", index = None, header=False, startrow=0)
                                        plot_df.to_excel(writer, sheet_name="pearson r and p-values", index = None, header=True, startrow= 8)
                                    
                                    print('')
                                    print('Correlation matrix saved')

                                    if valCorrelation['presentation_new'] == True or valCorrelation['presentation_old'] == True:
                                        filename = 'Correlation-table_' + var_tofilename + '_' + now +'.png'
                                        new_filename = os.path.join(output_folder, filename)
                                        df_save_to_image(plot, new_filename)                                        
                                        AddToPresentation(output_folder, new_filename)

                                if valCorrelation['presentation_new'] == True or valCorrelation['presentation_old'] == True:
                                    SavePresentation(valCorrelation)

                            break
                   
        #Regression

        if (not winRegression_active) and (valOriginal['Regression']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winRegression_active=True

            layoutRegression = [
                
                [sg.Text('Regression analysis:', justification='left', font=('Arial', 20))],
                [sg.Text('')],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('')],
                [sg.Text('Enter independent variable(s) (IVs) by name', font=('bold'))],
                [sg.Text('Variable names are defined by the Column headers (row 1) in the dataset.')],
                [sg.Text('For multiple regression, separate IVs by commas (no blank spaces).')],
                [sg.InputText('', size=(70,1), key='iv')],
                [sg.Text('Enter dependent variable (DV) and alternatively weight variable:.', font=('bold'))],
                [sg.InputText('', size=(50,1), key='dv'), sg.Text('Weight:'), sg.InputText('', size=(10,1), key='weight')],
                [sg.Text('Check to get standardized coeff. (Beta): '), sg.Radio('B', "RADIO1", key="std_no", default=True, size=(4,1)),  sg.Radio('Beta', "RADIO1", key="std_yes", default=False, size=(8,1))],
                [sg.Text('Export output: ', font=('bold'), size=(15,1)), sg.Radio('No: ', 'RADIO2', key='regression_output_no', default=True, size=(13,1)), sg.Radio('Yes: ', 'RADIO2', key='regression_output_yes', default=False, size=(13,1))],
                [sg.Text('Select output formats:', font=('bold'))],
                [sg.Checkbox('Regression plots', key='regression_plots', default=False, size=(13,1)), sg.Checkbox('Table', key='table', default=False, size=(13,1)), sg.Checkbox('VIF values (multiple regression only)', key='vif', default=False, size=(25,1))],
                
                [sg.Text('Select output folder:', font=('bold'))],
                [sg.InputText('', size=(70,1), key='regression_output'), sg.FolderBrowse()],
                [sg.Button('Continue'), sg.Button('Back')],
                [sg.Text('')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutRegression, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,600), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteRegression = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]
            
            winRegression=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteRegression)

            while True:
                evRegression, valRegression = winRegression.Read(timeout=100)               
                
                if evRegression=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evRegression==''

                if evRegression is None or evRegression=='Back' or evRegression=='Step back and change dataset':
                    winRegression_active=False
                    winRegression.Close()
                    del winRegression
                    winOriginal.UnHide()
                    break

                    
                if (evRegression=='Continue' and valRegression['iv']==''):
                    popup_break(evRegression, 'Please select independent variable(s)')

                elif (evRegression=='Continue' and valRegression['dv']=='' and valRegression['vif']==False):
                    popup_break(evRegression, 'Please select dependent variable')
                    
                elif (evRegression=='Continue') and (',' in valRegression['dv']):
                    popup_break(evRegression, 'Please select only one dependent variable.')

                elif (evRegression=='Continue') and (valRegression['regression_plots']==False) and (valRegression['regression_output_yes']==True) and (valRegression['table']==False) and (valRegression['vif']==False):
                    popup_break(evRegression, 'Please select output format')
                

                elif (evRegression=='Continue' and valRegression['regression_output']=='' and valRegression['regression_output_yes']==True):
                    popup_break(evRegression, 'Please select output folder')

                elif (evRegression=='Continue' and valRegression['regression_output_yes']==True and (valRegression['table']==False and valRegression['regression_plots']==False and valRegression['vif']==False)):
                    popup_break(evRegression, 'Please select output format(s)')
                #Valdidate variables supplied
                
                elif (evRegression=='Continue'):
                    iv=valRegression['iv']
                    iv=f'{iv}'
                    
                    #''.join(e for e in iv if e.isalnum())
                    #iv.strip() #Removes formatting and whitespace
                    
                    #iv=iv.replace(' ', '')
                    dv=valRegression['dv']
                    weight=valRegression['weight']
                    weight=f'{weight}'
                    
                    #Fix copy paste errors from variables overview:
                    if (weight!=""):
                        var=f'{iv},{dv},{weight}'
                    else:
                        var=f'{iv},{dv}'
                        
                    if var.endswith(','):
                        var=var.rstrip(',')
                    if var.endswith(', '):
                        var=var.rstrip(', ')
                    if var.startswith(','):
                        var=var.lstrip(',')
                    var=var.replace(', ', ',')
                    var=var.replace(' ','')
                    var=var.strip()
                    varsplit=var.split(',')

                    today = date.today()
                    now = today.strftime("%Y-%m-%d")

                    data=data_df[varsplit]
                    #Only numeric values: 
                    data = data.apply(pd.to_numeric, errors='coerce')
                    #Remove rows with non-integar data
                    #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                    #remove rows with empty cells
                    data = data.dropna(axis=0)
                    
                    if valRegression['std_yes']==True:
                        print('')
                        print('NB: Standardized values and coeffisient.')
                        print('')
                        if valRegression['weight'] !="":
                            data_main=data.iloc[: , :-1]
                            data_weight=data[data.columns[-1]]
                            # standardizing dataframe
                            #data = data.select_dtypes(include=[np.number]).dropna().apply(stats.zscore)
                            data_main_z = data_main.apply(zscore)
                            data=data_main_z.join(data_weight)

                        else:
                            data = data.apply(zscore)

                    
                    #Fjerne andre tegn på tomme celler - virker ikke
                    
                    validvariables=list(data.columns)
                    
                    list_var=list(varsplit)
                    
                    response='yes'
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'
 
                    
                    if (evRegression=='Continue') and (response =='no'):
                        popup_break(evRegression, 'Cancel. One or multiple variable(s) not in dataset')

                    #End validation of variables
                    while True:
                        
                        if (evRegression=='Continue') and (response=='yes'):
                            #pingouin
                    
                            iv=valRegression['iv']
                            iv.replace(", ", ",")
                            #iv_str="'"+iv+"'"
                            ivsplit=iv.split(',')#list of multiple IVs for multiple regression
                            dv=valRegression['dv']
                            output_folder=valRegression['regression_output']


                            if (',' in valRegression['iv']):
                                regression_type='multiple'
                            else:
                                regression_type='simple'

                            if regression_type=='simple':
                            
                                Y=data[dv]
                                X=data[iv]
                                
                                if weight == "":
                                    
                                    X=sm.add_constant(X)
                                    model=sm.OLS(Y, X).fit()
                                    predictions=model.predict(X)
                                    model.predict(X)

                                    print('Simple regression analysis: ')
                                    print('')
                                    print(model.summary())
                                
                                    print('Simple regression analysis completed')
                                    print('...')
                                    
                                else:
                                    weight_list = list(data[weight])
                                    X=sm.add_constant(X)
                                    model=sm.WLS(Y, X, weights=weight_list).fit()
                                    #predictions=model.predict(X)
                                    #model.predict(X)

                                    print('Simple weighted regression analysis: ')
                                    print('')
                                    print(model.summary())
                                
                                    print('Simple weighted regression analysis completed')
                                    print('...')

                                if valRegression['regression_output_yes']==True:
                                    iv_tofilename=iv.replace(',', '_')

                                    if valRegression['vif']==True:
                                        popup_break(evRegression, '"VIF - values" are only available for for multiple regression (multiple DVs)')
                    
                                    if valRegression['table']==True:#Make and save table:
                                        #model.summary()##Used with large amounts of variables.

                                        fig=plt.figure(figsize=(12, 7))
                                        plt.text(0.01, 0.05, str(model.summary()), {'fontsize': 10}, fontproperties = 'monospace')
                                        plt.axis('off')
                                        #plt.show()
                                        
                                        if (valRegression['std_yes'] == True) and (not valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_Beta_weighted_{now}_table.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_Beta_weighted_{now}_table.pdf'
                                        elif (valRegression['std_yes'] == False) and (not valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_B_weighted_{now}_table.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_B_weighted_{now}_table.pdf'
                                        elif (valRegression['std_yes'] == True) and (valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_Beta_{now}_table.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_Beta_{now}_table.pdf'
                                        else: 
                                            filename_png=f'Simple_OLS_{iv}_{dv}_B_{now}_table.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_B_{now}_table.pdf'
                                    
                                        plt.savefig(os.path.join(output_folder, filename_png), dpi=300, format='png', transparent=True)
                                        plt.savefig(os.path.join(output_folder, filename_pdf), dpi=300, format='pdf')
                                        plt.close(fig)
                                        print('Table saved')
                                        print('...')

                                    if valRegression['regression_plots']==True:

                                        fig, ax = plt.subplots(figsize=(15,9))
                                        plt.axis('off')
                                        fig=sm.graphics.plot_partregress_grid(model, fig=fig)
                                        #alternative: plot_ccpr_grid
                                        fig.tight_layout(pad=1.0)
                                        #plt.show()

                                        if (valRegression['std_yes'] == True) and (not valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_Beta_weighted_{now}_plot.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_Beta_weighted_{now}_plot.pdf'
                                        elif (valRegression['std_yes'] == False) and (not valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_B_weighted_{now}_plottable.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_B_weighted_{now}_plot.pdf'
                                        elif (valRegression['std_yes'] == True) and (valRegression['weight'] == ""):
                                            filename_png=f'Simple_OLS_{iv}_{dv}_Beta_{now}_plot.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_Beta_{now}_plot.pdf'
                                        else: 
                                            filename_png=f'Simple_OLS_{iv}_{dv}_B_{now}_plot.png'
                                            filename_pdf=f'Simple_OLS_{iv}_{dv}_B_{now}_plot.pdf'
                                        
                                        plt.savefig(os.path.join(output_folder, filename_png), dpi=300, format='png', transparent=True)
                                        plt.savefig(os.path.join(output_folder, filename_pdf), dpi=300, format='pdf')
                                        plt.close(fig)
                                        print('Plots saved as png and pdf')
                                        print('...')

                                #break    

                            if regression_type=='multiple':
  
                                Y=data[dv]
                                X=data[ivsplit]

                                if valRegression['vif']==True: #Make and save VIF values
                                                        
                                    if (evRegression=='Continue') and not (valRegression['weight']==""):
                                        popup_break(evRegression, 'VIF values only available for non-weighted data')
                                        
                                    vif = pd.DataFrame()
                                    vif["VIF Factor"] = [variance_inflation_factor(X.values, i) for i in range(X.shape[1])]
                                    vif["Variables"] = X.columns
                                    form=vif.round(1)

                                    df_vif=pd.DataFrame(form)
                                    print("VIF-values:")
                                    print('')
                                    print(df_vif)
                                    print('')

                                if weight == "" and not dv=='':
                                    X=sm.add_constant(X)
                                    model=sm.OLS(Y, X).fit()
                                    predictions=model.predict(X)
                                    model.predict(X)

                                    
                                    print('Multiple regression analysis: ')
                                    print('')
                                    print(model.summary())
                                    if valRegression['std_yes']==True:
                                        print("")
                                        rmse_residuals = np.sqrt(model.mse_resid)
                                        sd_str = str(rmse_residuals)
                                        print('Standard Deviation: ' + sd_str)
                                        print("")
                                        sd_list=['Standard deviation (sd): ', sd_str ]
                                        sd_df = pd.DataFrame(sd_list, index=None, columns=None)
                                    print('Multiple regression analysis completed')
                                    print('...')
                                    
                                elif not weight == "" and not dv=='':
                                    weight_list = list(data[weight])
                                    X=sm.add_constant(X)
                                    model=sm.WLS(Y, X, weights=weight_list).fit()
                                    #predictions=model.predict(X)
                                    #model.predict(X)

                                    print('Multiple weighted regression analysis: ')
                                    print('')
                                    print(model.summary())
                                    
                                    if valRegression['std_yes']==True:
                                        print("")
                                        rmse_residuals = np.sqrt(model.mse_resid)
                                        sd_str = str(rmse_residuals)
                                        print('Standard Deviation: ' + sd_str)
                                        print("")
                                        sd_list=['Standard deviation (sd): ', sd_str ]
                                        sd_df = pd.DataFrame(sd_list, index=None, columns=None)
                                
                                    print('Multiple weighted regression analysis completed')
                                    print('...')

                                if valRegression['regression_output_yes']==True:

                                    if valRegression['table']==True and not dv=='': #Save table (model summary):
                                        
                                        #print(model.summary())##Used with large amounts of variables.
                                        fig=plt.figure(figsize=(12,7))
                                        plt.text(0.01, 0.05, str(model.summary()), {'fontsize': 10}, fontproperties = 'monospace')
                                        plt.axis('off')
                                        #plt.show()

                                        #Filenames with all variables:
                                        #filename_png=f'Multiple_OLS_{iv}_{dv}_table.png'
                                        #filename_pdf=f'Multiple_OLS_{iv}_{dv}_table.pdf'
                                        #Filenames with dv + datetime: 
                                        if (valRegression['std_yes'] == True) and (not valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_Beta_weighted_{now}_table.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_Beta_weighted_{now}_table.pdf'
                                        elif (valRegression['std_yes'] == False) and (not valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_B_weighted_{now}_table.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_B_weighted_{now}_table.pdf'
                                        elif (valRegression['std_yes'] == True) and (valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_Beta_{now}_table.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_Beta_{now}_table.pdf'
                                        else: 
                                            filename_png=f'Multiple_OLS_{dv}_B_{now}_table.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_B_{now}_table.pdf'

                                        
                                        print(filename_png)
                                        print(filename_pdf)
                                        
                                        fig.savefig(os.path.join(output_folder, filename_png), dpi=300, format='png', transparent=True)
                                        fig.savefig(os.path.join(output_folder, filename_pdf), dpi=300, format='pdf')
                                        plt.close(fig)
                                        print('Table saved')
                                        print ('...')
                                        
                                    if valRegression['regression_plots']==True and not dv=='': #Make and save ccpr plots:
                                        fig, ax = plt.subplots(figsize=(15,9))
                                        plt.axis('off')
                                        fig=sm.graphics.plot_ccpr_grid(model, fig=fig)
                                        fig.tight_layout(pad=1.0)
                                        #plt.show()
                                        
                                        #Filenames with all variables:
                                        #filename_png=f'Multiple_OLS_{iv}_{dv}_table.png'
                                        #filename_pdf=f'Multiple_OLS_{iv}_{dv}_table.pdf'
                                        #Filenames with dv + datetime:
                                        
                                        if (valRegression['std_yes'] == True) and (not valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_Beta_weighted_{now}_plot.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_Beta_weighted_{now}_plot.pdf'
                                        elif (valRegression['std_yes'] == False) and (not valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_B_weighted_{now}_plot.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_B_weighted_{now}_plot.pdf'
                                        elif (valRegression['std_yes'] == True) and (valRegression['weight'] == ""):
                                            filename_png=f'Multiple_OLS_{dv}_Beta_{now}_plot.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_Beta_{now}_plot.pdf'
                                        else: 
                                            filename_png=f'Multiple_OLS_{dv}_B_{now}_plot.png'
                                            filename_pdf=f'Multiple_OLS_{dv}_B_{now}_plot.pdf'
                                        

                                        plt.savefig(os.path.join(output_folder, filename_png), dpi=300, format='png', transparent=True)
                                        plt.savefig(os.path.join(output_folder, filename_pdf), dpi=300, format='pdf')
                                        plt.close(fig)

                                        print('Plots saved')
                                        print('...')

                                    if valRegression['vif']==True: #Save VIF values
                                                        
                                        engine = 'xlsxwriter'
                                        
                                       
                                        if (valRegression['std_yes'] == True) and (valRegression['weight'] == ""):
                                            filename_xlsx=f'VIF_{dv}_Beta_{now}.xlsx'
                                        else: 
                                            filename_xlsx=f'VIF_{dv}_B_{now}.xlsx'

                
                                        
                                        #print(df)
                                        with pd.ExcelWriter(os.path.join(output_folder, filename_xlsx), engine=engine) as writer:
                                            df_vif.to_excel(writer, sheet_name="VIF values", index = None, header=True)

                                        print('VIF values saved')
                                        print('...')
                                
                        break

        #Mediation

        if not winMediation_active and valOriginal['Mediation']==True and evOriginal=='Continue' and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winMediation_active=True

            layoutOriginalFrame = [
                [sg.Text('')],
                [sg.Text("Mediation analysis", size=(25,1), justification='left', font=("Arial", 20, 'bold'))],
                #[sg.Text('Based on the Process Macro by Andrew F. Hayes, Ph.D. (www.afhayes.com)')],
                #[sg.Text("The following analysis is equv. to Hayes' Process, model 4:")],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('')],
                [sg.Text('Enter independent variable name (single):', font=('bold'))],

                [sg.InputText('', key='iv', size=(20,1))],

                [sg.Text('Enter mediating variable (single):', font=('bold'))],
                [sg.InputText('', key='m', size=(20,1))],
                [sg.Text('Enter covariate variable (single, optional):', font=('bold'))],
                [sg.InputText('', key='cov', size=(20,1))],

                [sg.Text('Enter dependent variable (single):', font=('bold'))],
                [sg.InputText('', key='dv', size=(20,1))],
                [sg.Text('Export output: ', font=('bold')), sg.Radio('No: ', 'RADIO2', key='mediation_output_no', default=True, size=(20,1)), sg.Radio('Yes: ', 'RADIO2', key='mediation_output_yes', default=False, size=(20,1))],
                [sg.Text('Select output folder:', size=(35, 1), font=('bold'))],      
                [sg.InputText('', key='mediation_output', size=(60,1)), sg.FolderBrowse(target='mediation_output')],           
                [sg.Text('')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutMediation = [
                [sg.Frame('Function', layoutOriginalFrame, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,600), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteMediation = [[sg.Column(layoutMediation, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winMediation=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteMediation)


            while True:
                evMediation, valMediation = winMediation.Read(timeout=100)
                response='none'


                if evMediation=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ', '
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evMediation==''
                

                #Exceptions:


                if (evMediation is None) or (evMediation == 'Back') or (evMediation=='Step back and change dataset'):
                    winMediation_active=False
                    winMediation.Close()
                    del winMediation
                    winOriginal.UnHide()
                    break

                elif (evMediation=='Continue' and valMediation['iv']==''):
                    popup_break(evMediation, 'Please select independent variable')

                elif (evMediation=='Continue') and (',' in valMediation['iv']):
                    popup_break(evMediation, 'Please select only one independent variable')
                        
                elif (evMediation=='Continue' and valMediation['m']==''):
                    popup_break(evMediation, 'Please select mediating variable')

                elif (evMediation=='Continue' and valMediation['dv']==''):
                    popup_break(evMediation, 'Please select dependent variable')
                        
                elif (evMediation=='Continue') and (',' in valMediation['dv']):
                    popup_break(evMediation, 'Please select only one mediating variable')

                elif (evMediation=='Continue' and (valMediation['mediation_output_yes']==True) and valMediation['mediation_output']==''):
                    popup_break(evMediation, 'Please select output folder')
                
                #Valdidate variables supplied
                elif (evMediation=='Continue'):
                    
                    med_var_input=valMediation['m']

                    if ',' in med_var_input:
                        med_var = ''
                        med_var_list=med_var_input.split(',')

                        for var in med_var_list:
                            if med_var == '':
                                med_var=var

                            elif med_var != '':
                                med_var=med_var + ',' + var
    
                       
                    if ',' not in med_var_input:
                        med_var=med_var_input

                    var=valMediation['iv'] + ',' + med_var + ',' + valMediation['dv']
                    var.replace(" ", "")
                    if not valMediation['cov']=='':
                        cov = valMediation['cov']
                        var=var + ',' + cov

                    varsplit=var.split(',')
                    
                    validvariables=list(data_df.columns)
                   
                    list_var=list(varsplit)
                    
                    response='yes'
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'

                    if (evMediation=='Continue') and (response =='no'):
                        response='none'
                        popup_break(evMediation, 'One or several entered variables not in dataset')
                    
                    elif (evMediation=='Continue') and (response=='yes'):
                        
                        data=data_df[varsplit]
                        #Use numeric values only
                        data = data.apply(pd.to_numeric, errors='coerce')
                        #remove rows with strings and floats:
                        #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                        #remove rows with empty cells
                        data = data.dropna(axis=0)
                        

                        iv=valMediation['iv']
                        dv=valMediation['dv']
                       
                        intro_1 = 'MEDIATION ANALYSIS'
                        intro_2 = ''
                        intro_3 = 'Independent variable: ' + iv
                        intro_4 = 'Mediating variable(s): ' + med_var_input
                        intro_4_1 = 'Covariate(s). ' + valMediation['cov']
                        intro_5 = 'Dependent variable: ' + dv
                        intro_6 = 'Confidence interval: 95%'
                        intro_7 = 'Bootstrap number: 5000'
                        intro_8 = ''
                        intro_9 = 'Results:'

                        intro_1_list=list(intro_1.split(';'))
                        intro_2_list=list(intro_2.split(';'))
                        intro_3_list=list(intro_3.split(';'))
                        intro_4_list=list(intro_4.split(';'))
                        intro_4_1_list=list(intro_4_1.split(';'))
                        intro_5_list=list(intro_5.split(';'))
                        intro_6_list=list(intro_6.split(';'))
                        intro_7_list=list(intro_7.split(';'))
                        intro_8_list=list(intro_8.split(';'))
                        intro_9_list=list(intro_9.split(';'))

                        #variables_df=pd.DataFrame(variables_list)                    
                        print(intro_1)
                        print(intro_2)
                        print(intro_3)
                        print(intro_4)
                        print(intro_4_1)
                        print(intro_5)
                        print(intro_6)
                        print(intro_7)
                        print(intro_8)
                        print(intro_9)


                        if ',' not in med_var_input:
                            m=med_var_input
                            if(valMediation['cov'] == ''):
                                med_results = mediation_analysis(data=data, x=iv, m=m, y=dv, covar=None, alpha=0.05, n_boot= 5000, seed=42, return_dist=False)                            
                            else:
                                med_results = mediation_analysis(data=data, x=iv, m=m, y=dv, covar=cov, alpha=0.05, n_boot= 5000, seed=42, return_dist=False)

                            print(med_results.to_string(index=False))
                            print('Mediation analysis executed')
                            print('...')

                        else:
                            m=med_var_input
                            m=m.split(',')
                            if(valMediation['cov'] == ''):
                                med_results = mediation_analysis(data=data, x=iv, m=m, y=dv, covar=None, alpha=0.05, n_boot=5000, seed=None).round(3)
                            else:
                                med_results = mediation_analysis(data=data, x=iv, m=m, y=dv, covar=cov, alpha=0.05, n_boot=5000, seed=None).round(3)

                            print(med_results.to_string(index=False))
                            print('Mediation analysis executed')
                            print('...')

                        #End mediation by pingouin

                        engine = 'xlsxwriter'

                        if valMediation['mediation_output_yes']==True:
                            output_folder=valMediation['mediation_output']         

                            new_file_name='Mediation_' +  var + '.xlsx'              

                            with pd.ExcelWriter(os.path.join(output_folder, new_file_name), engine=engine) as writer:
                                med_results.to_excel(writer, sheet_name="Model_summary", index = None, startrow=11, startcol=0, header=True)
                                pd.DataFrame(intro_1_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=0, startcol=0, header=False)
                                pd.DataFrame(intro_2_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=1, startcol=0, header=False)
                                pd.DataFrame(intro_3_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=2, startcol=0, header=False)
                                pd.DataFrame(intro_4_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=3, startcol=0, header=False)
                                pd.DataFrame(intro_4_1_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=4, startcol=0, header=False)
                                pd.DataFrame(intro_5_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=5, startcol=0, header=False)
                                pd.DataFrame(intro_6_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=6, startcol=0, header=False)
                                pd.DataFrame(intro_7_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=7, startcol=0, header=False)
                                pd.DataFrame(intro_8_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=8, startcol=0, header=False)
                                pd.DataFrame(intro_9_list).to_excel(writer, sheet_name="Model_summary", index = None, startrow=9, startcol=0, header=False)

                            print('')
                            print ('Output saved to file')   

                        print('')
                        print('Done.')
                        print('')
                        print('- - -')
                        print('')       
                
        #Moderation

        if not winModeration_active and valOriginal['Moderation']==True and evOriginal=='Continue' and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winModerarion_active=True

            layoutOriginalFrame = [
                [sg.Text("Moderation analysis", size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Based on the Process Macro by Andrew F. Hayes, Ph.D. (www.afhayes.com)')],
                [sg.Text("The following analysis is equv. to Hayes' Process, model 1:")],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('Use the column headers (row 1) in the dataset as variable names.')],
                [sg.Text('Enter independent variable (IV):', font=('bold'))],
                [sg.InputText('', key='iv', size=(20,1))],
                [sg.Text('Enter moderating variable (M): (NB: Has to be categorical.)', font=('bold'))],
                [sg.InputText('', key='m', size=(20,1))],
                [sg.Text('Enter dependent variable (DV):', font=('bold'))],
                [sg.InputText('', key='dv', size=(20,1))],
                [sg.Text('Export output: ', font=('bold')), sg.Radio('No: ', 'RADIO2', key='moderation_output_no', default=True, size=(20,1)), sg.Radio('Yes: ', 'RADIO2', key='moderation_output_yes', default=False, size=(20,1))],
                [sg.Text('Choose output format:', font=('bold'))],
                [sg.Checkbox('Save line chart as pdf', key='pdf', default=False, size=(22,1)), sg.Checkbox('Save line chart as png', key='png', default=False, size=(22,1)), sg.Checkbox('Export values to Excel', key='xlsx', default=False, size=(22,1))],    
                [sg.Text('Select output folder:', size=(35, 1), font=('bold'))],      
                [sg.InputText('', key='moderation_output', size=(60,1)), sg.FolderBrowse(target='moderation_output')],           
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutModeration = [
                [sg.Frame('Function', layoutOriginalFrame, size= (570, 600))]]
            
            layoutRight = [
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,600), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteModeration = [[sg.Column(layoutModeration, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winModeration=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteModeration)

            while True:
                evModeration, valModeration = winModeration.Read(timeout=100)
                response='none'

                if evModeration=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ', '
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evModeration==''
 
                #Exceptions:

                if evModeration is None or evModeration == 'Back' or evModeration=='Step back and change dataset':
                    winModeration_active=False
                    winModeration.Close()
                    del winModeration
                    winOriginal.UnHide()
                    break
                      
                if (evModeration=='Continue' and valModeration['iv']==''):
                    popup_break(evModeration, 'Please select independent variable')
                    
                elif (evModeration=='Continue' and ',' in valModeration['iv']):
                    popup_break(evModeration, 'Please select only one independent variable')
                        
                elif (evModeration=='Continue' and valModeration['m']==''):
                    popup_break(evModeration, 'Please select moderating variable')

                elif (evModeration=='Continue' and ',' in valModeration['m']):
                    popup_break(evModeration, 'Please select only one moderating variable')
                    
                elif (evModeration=='Continue' and valModeration['dv']==''):
                    popup_break(evModeration, 'Please select dependent variable')
                    
                elif (evModeration=='Continue' and ',' in valModeration['dv']):
                    popup_break(evModeration, 'Please select only one dependent variable')
                
                elif (evModeration=='Continue' and valModeration['moderation_output_yes']==True and valModeration['moderation_output']==''):
                    popup_break(evModeration, 'Please select output folder')


                 #Valdidate variables supplied
                elif (evModeration=='Continue'):

                    iv=valModeration['iv']
                    dv=valModeration['dv']
                    mod=valModeration['m']
                    med_list=list(mod)
                    output_folder=valModeration['moderation_output']
                    validvariables=list(data_df.columns)
                    var=valModeration['iv']+','+valModeration['m']+','+valModeration['dv']
                    var=var.replace(" ", "")
                    varsplit=var.split(',')
                    list_var=list(varsplit)

                    #response='yes'
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'
                            
                        else:
                            response='yes'

                    if response=='no':
                        response='none'
                        popup_break(evModeration, 'One or several entered variables not in dataset')

                    data=data_df[varsplit]
                    #Use numeric values only
                    data = data.apply(pd.to_numeric, errors='coerce')
                    #remove rows with strings and floats:
                    #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                    data = data.dropna(axis=0)
                    

                    moderatorvalues=[]
                    
                    for i in data[mod]:
                        mod_str=str(i)
                        if mod_str not in moderatorvalues:
                            moderatorvalues.append(mod_str)

                    #check if moderator is dictonomous
                    if not len(moderatorvalues)==2:
                        popup_break(evModeration, 'Moderating variable is not dictonomous (only two values)')
                    
                                        
                    if (evModeration=='Continue') and (response=='yes'):
                        
                        
                    
                        p = Process(data=data, model=4, x=iv, y=dv, m=[mod])
                        summary = p.summary()
                        model_modskills=p.outcome_models

                        direct_model=p.direct_model
                        direct_summary=direct_model.coeff_summary()
                        
                        dirEff=direct_summary.iloc[0,0]
                        
                        indirect_summary=p.indirect_model.coeff_summary()
                        indEff=indirect_summary.iloc[0,1]
                        
                        
                        #Getting constant through OLS
                    
                        reg = linear_model.LinearRegression()
                        i_var=valModeration['iv']+','+valModeration['m']
                        i_varsplit=i_var.split(',')
                        
                        X=data[i_varsplit]              
                        Y=data[dv]
                                
                        X=sm.add_constant(X)
                        cons_df=pd.DataFrame(X)
                        
                        model=sm.OLS(Y, X).fit()     
                        params=model.params
                        const=(params[0])
                        
                        redictions=model.predict(X)
                        sum=model.summary()

                        fig = plt.figure(figsize=(12,8))

                        plt.title('Moderation analysis: IV = ' + iv + ', M = ' + mod + ', DV = ' + dv)
                        #x=np.linspace(1,5,num=5)

                        if const >= (const + dirEff) and const >=(const + indEff):
                            upperLimit=const + 0.25
                        elif (const + dirEff) >= (const + indEff):
                            upperLimit=const + dirEff + 0.25
                        else:
                            upperLimit = const + indEff +0.25

                        if const <= (const + dirEff) and const <=(const + indEff):
                            lowerLimit=const -0.25
                        elif (const + dirEff) <= (const + indEff):
                            lowerLimit=const + dirEff - 0.25
                        else:
                            lowerLimit = const + indEff - 0.25

                        xDir, yDir = [0, 1], [const, const + dirEff]
                        xInd, yInd = [0, 1], [const, const + indEff]
                        
                        ax = plt.subplot(111)
                        plt.xlabel('Level of '+iv)
                        ax.spines['right'].set_visible(False)
                        ax.spines['top'].set_visible(False)
                        plt.ylabel('Level of '+dv)
                        
                        plt.plot(xDir, yDir, 'b', label='Direct effect ', linewidth=3)
                        plt.plot(xInd, yInd, 'r', label='Indirect effect', linewidth=3)
                        plt.xlim(0,1), plt.ylim(lowerLimit,upperLimit)
                        plt.legend()
                                         

                        if valModeration['moderation_output_yes'] == True:
                            
                            
                            if valModeration['moderation_output_yes'] == True:

                                if valModeration['xlsx']==True:
                            
                                    header=pd.DataFrame(columns=["Results from the moderation analysis:"])
                                    header_1=pd.DataFrame(columns=['IV: ' + str(iv) + '    DV: ' + str(dv) + '    Moderator: ' + str(mod)])
                                    header_2=pd.DataFrame(columns=['Direct effect of ' + str(iv) + ' on ' + str(dv) + ':'])
                                    header_3=pd.DataFrame(columns=['Indirect effect of ' + str(iv) + ' on ' + str(dv) + ', moderatad by ' + str(mod) + ':'])
                                    engine = 'xlsxwriter'
                                    filename='Moderation_'+iv+'_'+ mod +'_'+dv+'.xlsx'
                                    new_file=os.path.join(output_folder, filename)
                                    new_file=os.path.normpath(new_file)


                                    with pd.ExcelWriter(new_file, engine=engine) as writer:
                                        header.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=1, startrow=0)
                                        header_1.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=1, startrow=1)
                                        header_2.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=1, startrow=2)
                                        direct_summary.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=-1, startrow=5, header=True)
                                        header_3.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=2, startrow=8)
                                        indirect_summary.style.set_properties({'text-align': 'left'}).to_excel(writer, sheet_name="Results", startcol=-2, startrow=9, header=True)

                                    print('Values saved to Excel')
                                    print('...')
                                

                                if valModeration['png']==True:
                                    plt.savefig(os.path.join(output_folder, 'Moderation_')+iv+'_'+mod+'_'+dv+'_moderation.png', dpi=300, format='png', transparent=True)
                                    print('Line chart saved as image file')
                                    print('...')
                                if valModeration['pdf']==True:
                                    plt.savefig(os.path.join(output_folder, 'Moderation_')+iv+'_'+mod+'_'+dv+'_moderation.pdf', dpi=300, format='pdf', transparent=True)
                                    print('Line chart saved as PDF file')
                                    print('...')

                                plt.close(fig)
                #break

        
        #Compute variables

        if (not winCompute_active) and (valOriginal['Compute']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winCompute_active=True


            layoutCompute = [
                [sg.Text('Edit dataset:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('Compute new variable', font=('bold'))],
                [sg.Text('Enter new variable name (NB:No white space or special characters):', size=(50,1)), sg.InputText('', key='new_variable', size=(15, 1))],
                [sg.Text('Choose variable category: ', font=('bold')), sg.Radio('Median (M): ', 'RADIO1', key='variable_average', default=False, size=(10,1)), sg.Radio('Dictonomous: ', 'RADIO1', key='variable_dict', default=False, size=(20,1))],
                [sg.Text('For Mean (M), enter source variables, divided by comma.')],
                [sg.Text('For dictonomous enter source variable and selected value (integer), divided by comma.')],
                [sg.InputText('', key='compute_variables', size=(35,1))],  
                [sg.Text('Delete variable:', font=('bold'))],
                [sg.Text('Enter variable to be deleted:', size=(25,1)), sg.InputText('', key='del_variables', size=(15, 1))],
                [sg.Text('NB: To delete multiple variables, enter variables separated by comma(,)')],
                [sg.Text('Select output folder:', font=('bold'), size=(60, 1))],      
                [sg.InputText('', key='compute_output'), sg.FolderBrowse()],
                [sg.Text('NB: If left open, source folder will be used')],
                [sg.Text('Enter new Excel file name (excl. .xlsx):', font=('bold'), size=(35, 1))],      
                [sg.InputText('', key='compute_filename', size=(35, 1)), sg.Text('NB: If left open, original filename will be used')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutCompute, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('This tool creates new variable in the dataset', font=('bold'), size = (50, 1), justification='left', background_color='lightsteelblue2')], 
                [sg.Text('Mean (M) creates a new variable with the mean of selected variables.', size=(50,1), background_color='lightsteelblue2')],
                [sg.Text('Dictonomous creates a categorical variable based on a source variable:', size=(55,1), background_color='lightsteelblue2')],
                [sg.Text('Observations with selected value are coded to 1, other values 0 and missing to "NaN".', size=(50,1), background_color='lightsteelblue2')]], size=(500,150), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(60,50))]], size=(505,450), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteCompute = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winCompute=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteCompute)

            while True:
                evCompute, valCompute = winCompute.Read(timeout=100)    

                if evCompute=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evCompute==''

                if evCompute is None or evCompute=='Back' or evCompute=='Step back and change dataset':
                    winCompute_active=False
                    winCompute.Close()
                    del winCompute
                    winOriginal.UnHide()
                    break       

                now = datetime.now()
                now=now.strftime('%d/%m/%Y %H:%M:%S')
                now=str(now)

                if (evCompute=='Continue' and (valCompute['compute_variables']=='' and valCompute['del_variables']=='')):
                    popup_break(evCompute, 'Choose variables')
                elif (evCompute=='Continue' and (valCompute['new_variable']=='' and not valCompute['compute_variables']=='')):
                    popup_break(evCompute, 'Choose new variable name')
                elif (evCompute=='Continue' and valCompute['compute_variables']=='' and valCompute['del_variables']==''):
                    popup_break(evCompute, 'Choose variables')

                #Valdidate variables supplied
                elif (evCompute=='Continue'):
                    validvariables=list(data_df.columns)
                    new=valCompute['new_variable']
                    if new in validvariables:
                        popup_break(evCompute, 'New variable already exists')
                    
                    response='yes'

                    if valCompute['variable_average']==True:
                        var=valCompute['compute_variables']
                        var=var.replace(" ", "")
                        varsplit=var.split(",")

                        for string in varsplit:    
                            if string not in validvariables:
                                response='no'     

                    elif valCompute['variable_dict']==True:
                        content=valCompute['compute_variables']
                        content=content.replace(" ", "")
                        contsplit=content.split(',')
                        list_cont=list(contsplit)
                        len_cont=len(list_cont)
                        
                        if not len_cont == 2:
                            popup_break(evCompute, 'Please select only one variable and one value. Selected value must be an integer or string.')
                        
                        dict_value=int(list_cont[-1])
                        var=str(list_cont[0])
                        
                        if var not in validvariables:
                            response='no'

                    #Checking existing sheets
                    dataset_sheets=pd.ExcelFile(data_in_str)
                    sheets_list=dataset_sheets.sheet_names
                    sheets_str=' '.join(sheets_list)

                    print("Existing sheets: " + sheets_str)

                    if "Changed variables syntax" in sheets_list:
                        print('Syntax sheet already exists')
                        syntax_df=dataset_sheets.parse("Changed variables syntax")
                    
                    else:
                        #Create syntax sheet
                        syntax_df=pd.DataFrame(columns=['Timestamp', 'Variable', 'Details'])
                        print('New syntax sheet ready for content')

                    while True:
            
                        if (evCompute=='Continue') and (response =='no'):
                            popup_break(evCompute, 'Selected variable not in dataset')    

                        if valCompute['compute_filename'] == '':
                            new_filename=chosen_filename
                        else:
                            new_filename=valCompute['compute_filename'] + '.xlsx'

                        if valCompute['compute_output'] == '':
                            out_folder=os.path.dirname(data_in)
                        else:
                            out_folder=valCompute['compute_output'] 

                        if valCompute['variable_average'] == True and (response=='yes'):

                            print('Creating new variable: ' + new)
                            data_df[new] = data_df[varsplit].mean(axis=1)
                            print('New variable added to dataset')
                            print('')
                            syntax_new_dict={'Timestamp': now, 'Variable': new, 'Details': 'DataFrame' + str(varsplit) + '.mean(axis=1)'}
                            print('Details: ' + str(syntax_new_dict))
                            syntax_df.loc[len(syntax_df.index)] = syntax_new_dict
                            new_syntax_df=syntax_df
                            print('Details added to syntax sheet')  
                            print(new_syntax_df)
                            print('')

                        elif valCompute['variable_dict']==True and (response=='yes'):
                            for i in data_df.index:
                                #if (isinstance(data_df.at[i, var], int)):
                                if data_df.at[i, var]==dict_value: 
                                    data_df.at[i, new] = 1
                                elif data_df.at[i, var]=='': 
                                    data_df.at[i, new] = ''
                                elif data_df.at[i, var]=='NaN': 
                                    data_df.at[i, new] = 'NaN'    
                                else:
                                    data_df.at[i, new] = 0
                            
                            print('Creating new variable: ' + new)
                            print('New variable added to dataset')
                            details_syntax=str('Dict: ' + str(var) + ': ' + str(dict_value) + ' = 1, else = 0')
                            print('Details: ' + details_syntax)
                            #syntax_new_dict={'Timestamp': now, 'Variable': new, 'Details': details_syntax}
                            syntax_new_list=[now, new, details_syntax]
                            #syntax_df.loc[len(syntax_df)] = syntax_new_dict
                            syntax_df.loc[len(syntax_df.index)] = syntax_new_list
                            print('Details added to syntax sheet')
                            new_syntax_df=syntax_df
                            print('')
                        
                        elif evCompute=='Continue' and not valCompute['del_variables']=='':
                            response='yes'
                            
                            var=valCompute['del_variables']
                            var=var.replace(" ", "")
                            var=str(var)
                            varsplit=var.split(",")
                            len_var=len(varsplit)

                            if len_var > 1:
                                del_var=varsplit
                                print('Variables to be deleted: ' + var)
                            else:
                                del_var=var
                                print('Variable to be deleted: ' + var)
                            
                            for string in varsplit:    
                                if string not in validvariables:
                                    response='no'
                            
                            if response =='no':
                                popup_break(evCompute, 'Selected variable(s) not in dataset')   

                            data_df=data_df.drop(del_var, axis=1)
                            print('Variable(s) deleted')
                            
                            if valCompute['compute_filename'] == '':
                                new_filename=chosen_filename
                            else:
                                new_filename=valCompute['compute_filename'] + '.xlsx'

                            if valCompute['compute_output'] == '':
                                out_folder=os.path.dirname(data_in)
                            else:
                                out_folder=valCompute['compute_output'] 
                            
                            syntax_new_list=[now, var, 'Variable(s) deleted']
                            syntax_df.loc[len(syntax_df.index)] = syntax_new_list
                            new_syntax_df=syntax_df

                        engine = 'xlsxwriter'
                        new_excel_file=os.path.join(out_folder, new_filename)
                        new_excel_file=os.path.normpath(new_excel_file)

                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                            data_df.to_excel(writer, sheet_name="Edited dataset", index = None, header=True, startrow=0, startcol=0)
                            new_syntax_df.to_excel(writer, sheet_name="Changed variables syntax", index = None, header=True, startrow=0, startcol=0)
                        print('')
                        print(' New dataset saved')
                        print('')

                        data_in=data_in_str
                        
                        break

        #Factor analysis
        
        if (not winFactor_active) and (valOriginal['Factor']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winFactor_active=True

            layoutFactor = [
                [sg.Text('Factor analysis:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('')],
                [sg.Text('Choose action:', font=('bold'))],
                [sg.Radio("Bartlett's test: ", 'RADIO1', key='bartlett', default=False, size=(17,1)), sg.Radio('Kaiser-Mayer-Olkin test: ', 'RADIO1', key='kmo', default=False, size=(17,1)), sg.Radio('Get factors and eigenvalues: ', 'RADIO1', key='kai_crit', default=False, size=(21,1))], 
                [sg.Radio('Exploratory analysis: ', 'RADIO1', key='exploratory', default=False, size=(17,1)), sg.Radio('Confirmatory analysis: ', 'RADIO1', key='confirmatory', default=False, size=(17,1))],
                [sg.Text('For tests or analyses:', font=('bold'), size=(35,1))],
                [sg.Text('Variables:', size=(13,1)), sg.InputText('', key='factor_variables', size=(52,3))],
                [sg.Text('For exploratory analyses:', size=(20,1)),sg.Text('Number of factors: ', size=(13, 1)),sg.InputText('', key='nr_factors', size=(3, 1))],
                [sg.Text('Rotation: ', size=(13,1)), sg.Radio('Varimax:', 'RADIO2', key='varimax', default=False, size=(6,1)), sg.Radio('Promax:', 'RADIO2', key='promax', default=False, size=(6,1)), sg.Radio('Oblimin:', 'RADIO2', key='oblimin', default=False, size=(6,1)), sg.Radio('Oblimax:', 'RADIO2', key='oblimax', default=False, size=(6,1)), sg.Radio('None:', 'RADIO2', key='none', default=True, size=(6,1))],
                [sg.Text('Output (optional):', font=('bold'), size=(35, 1))],
                [sg.Checkbox('Export table as .xlsx', key='factor_table', default=False, size=(17,1)), sg.Checkbox('Export scree plot as pdf', key='factor_scree_pdf', default=False, size=(17,1)), sg.Checkbox('Export scree plot as png', key='factor_scree_png', default=False, size=(19,1))],
                [sg.Text('Output filename:', size=(13,1)), sg.InputText('', key='factor_filename', size=(32,1)), sg.Text('(NB: Excl. file extension)', size=(20,1))],
                [sg.Text('Output folder:', size=(13,1)), sg.InputText('', key='factor_output', size=(52,1)), sg.FolderBrowse()],
                [sg.Text('', size=(13,1)), sg.Text('NB: If left open, dataset folder will be used')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutFactor, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('For all tests and analyses:', font=('bold'), size=(50,1), background_color='lightsteelblue2')],
                [sg.Text('Please enter either all included variables separated by commas, or enter excluded variables with a minus sign in front of the first variable. If left open, all variables in the dataset will be included.', size=(50,3), background_color='lightsteelblue2')],
                [sg.Text('For Confirmatory analysis, please enter factor names and variables in this format: Facor_1: var_1 var 2 Factor_2: var_1 var_2 etc ... .', size=(50,2), background_color='lightsteelblue2')],
                [sg.Text('For Confirmatory analysis: ', size=(55,1), font=('bold'), background_color='lightsteelblue2')]], size=(500,170), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(60,80))]], size=(505,430), background_color='lightsteelblue2')]]
            
            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteFactor = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]


            winFactor=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteFactor)

            while True:
                evFactor, valFactor = winFactor.Read(timeout=100)

                if evFactor=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evFactor==''
                
                if evFactor is None or evFactor=='Back' or evFactor=='Step back and change dataset':
                    winFactor_active=False
                    winFactor.Close()
                    del winFactor
                    winOriginal.UnHide()
                    break       

                #Valdidate variables supplied

                if(evFactor=='Continue' and valFactor['bartlett']==False and valFactor['kmo']==False and valFactor['kai_crit']==False and valFactor['exploratory']==False and valFactor['confirmatory']==False):
                    popup_break(evFactor, 'Please choose action.')
                elif(evFactor=='Continue' and (valFactor['factor_table']==True or valFactor['factor_scree_pdf']==True or valFactor['factor_scree_png']==True) and valFactor['factor_filename']==''):
                    popup_break(evFactor, 'Please enter a name for the export file.')
                elif (evFactor=='Continue' and (valFactor['exploratory']==True or valFactor['bartlett']==True or valFactor['kmo']==True or valFactor['confirmatory']==True) and (valFactor['factor_scree_png']==True or valFactor['factor_scree_pdf']==True)):
                    valFactor['factor_scree_png']=False
                    valFactor['factor_scree_pdf']=False
                    popup_break(evFactor, 'Scree plots are only available for the "Get factors and eigenvalues" option')
                elif(evFactor=='Continue' and valFactor['exploratory']==True and valFactor['nr_factors']==''):
                    #valFactor['confirmatory']=False
                    popup_break(evFactor, 'For exploratory analysis, please enter number of factors.')
                elif(evFactor=='Continue' and valFactor['confirmatory']==True and (':' not in valFactor['factor_variables'] or ',' in valFactor['factor_variables'])):
                    popup_break(evFactor, 'Please enter factors and variables in the right format.') 
                          
                elif evFactor=='Continue':
                    #evFactor=False
                    
                    while True:      
                        validvariables=list(data_df.columns)
                        filename=valFactor['factor_filename'] 

                        if valFactor['factor_variables'] == '':
                            varsplit=validvariables
                            varsplit=validvariables.replace(', ', ',')
                            data=data_df
                            #convert to numeric values
                            data = data.apply(pd.to_numeric, errors='coerce')
                            #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                            data = data.dropna(axis=0)

                        elif ':' in valFactor['factor_variables']:
                            factor_str = valFactor['factor_variables']

                            #factor_str = f'{factor_str}'
                            print(factor_str)
                            print('Exploratory analysis ...')
                            #factor_dumps=json.dumps(factor_str)
                            #print(factor_dumps)
                            #print(type(factor_dumps))
                            #factor_string = factor_str.replace("'", "\"")
                            #factor_dict = json.load(factor_string)
                            #factor_dict=json.load(factor_str)
                            #print(factor_dict)
                            #print(type(factor_dict))
                            
                            
                            input_cfa = valFactor['factor_variables']
                            factor_dict={}
                            n=1
                            factors=[]
                            input_list=input_cfa.split()
                            var = []
                            print(input_list)

                            for i in input_list:
                                if i.endswith(':'):
                                    i=i.replace(':', '')
                                    factors.append(i)
                                    fact_nr_str = str(n-1)
                                    factor_dict[fact_nr_str] = var
                                    var = []
                                    n += 1
                                else:
                                    variable=i
                                    var.append(variable)

                            fact_nr_str=str(n-1)
                            factor_dict[fact_nr_str] = var
                            (k := next(iter(factor_dict)), factor_dict.pop(k))

                            #print(object.keys(factor_dict))
                            column_names_list = factors #list(factor_dict.keys())
                            print(type(column_names_list))

                            cfa_dict={}
                            cfa_variables=valFactor['factor_variables']
                            cfa_variables_list=cfa_variables.split()
                            print(cfa_variables) 
                            print('Factor names:')
                            print(column_names_list)

                            n_fact=len(column_names_list)
                            print(f'Number of factors: {n_fact}')

                            variables_included_list = list(factor_dict.values())
                            variables_included_str = str(variables_included_list)
                            print(type(variables_included_str))
                            print(f'Variables_list: {variables_included_str}')
                            variables_included_str = f'{variables_included_str}'
                            variables_included_str=variables_included_str.replace("[[", "")
                            variables_included_str=variables_included_str.replace("]]", "")
                            variables_included_str=variables_included_str.replace("]", "")
                            variables_included_str=variables_included_str.replace("[", "")
                            variables_included_str=variables_included_str.replace("'", "")
                            variables_included_str=variables_included_str.replace('],[',',')
                            variables_included_str=variables_included_str.replace(", ", ",")
                            variables_included_str=variables_included_str.replace("] ,[", ",")
                            variables_included_str=variables_included_str.replace(" ", "")
                            print(f'Values included: {variables_included_str}')
                            variables_included_list_new=variables_included_str.split(',')

                            
                            varsplit=variables_included_str.split(',')
                            #data=data_df.drop(varsplit, axis=1)
                            data=data_df[varsplit].copy()

                            #convert to numeric values
                            data = data.apply(pd.to_numeric, errors='coerce')
                            #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                            data = data.dropna(axis=0)
                            print(data)
                            
                        
                        else:
                            var_=valFactor['factor_variables']
                            if var_.startswith('-'):
                                var=var_[1:]
                                if var.endswith(','):
                                    var=var.rstrip(',')
                                if var.startswith(','):
                                    var=var.lstrip(',')
                                var=var.replace(', ', ',')
                                var=var.replace(' ','')
                                var=var.strip()
                                varsplit=var.split(',')
                                var=str(var)
                                varsplit=var.split(',')
                                data=data_df.drop(varsplit, axis=1)
                                #convet to numeric values
                                data = data.apply(pd.to_numeric, errors='coerce')
                                #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                                data = data.dropna(axis=0)

                            else:
                                var=var_
                                if var.endswith(','):
                                    var=var.rstrip(',')
                                if var.startswith(','):
                                    var=var.lstrip(',')
                                var=var.replace(', ', ',')
                                var=var.replace(' ','')
                                var=var.strip()
                                varsplit=var.split(',')
                                #var=str(var)
                                #varsplit=var.split(",")
                                data=data_df[varsplit]
                                #convet to numeric values
                                data = data.apply(pd.to_numeric, errors='coerce')
                                #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                                data = data.dropna(axis=0)


                        #Convert to numeric
                        ##data=data.apply(pd.to_numeric, downcast='integer', errors='coerce')
                        #Use numeric values only
                        data = data.apply(pd.to_numeric, errors='coerce')
                        #remove rows with missing values:
                        data = data.dropna(axis=0)
                        #Converts floats to ints
                        ##data = data.astype(int)
                        #pd.set_option('display.max_rows', None)
                        fact_variables=list(data.columns)
                        str_variables='Included varables: '
                        n=0
                        for str in fact_variables:
                            if n==0:
                                str_variables=str_variables + str
                                n=1
                            else:
                                str_variables=str_variables + ', ' + str
                        
                        nr_variables=len(fact_variables)

                        response='yes'

                        for string in varsplit:    
                            if string not in validvariables:
                                response='no'
                        
                        if response=='no':
                            popup_break(evFactor, 'Error: One or more entered variables not in dataset')

                        #while True:
                        
                        if response=='yes':
                    
                            if valFactor['factor_output'] == '':
                                out_folder=os.path.dirname(data_in)
                            else:
                                out_folder=valFactor['factor_output'] 
                            
                                
                            fa=FactorAnalyzer(rotation=None)
                            factor_summary=fa.fit(data)
                            loadings=fa.loadings_
                            fa.fit(data)
                            ev, v = fa.get_eigenvalues()
                            #Bartlett's test
                            
                            if valFactor['bartlett']==True:
                                print("Bartlett's test:")

                                chi_square_value,p_value=calculate_bartlett_sphericity(data)

                                chi=chi_square_value
                                p_val=p_value
                                bartlett_info=[]
                                bartlett_info.append('Test: Bartlett')
                                bartlett_info.append('')
                                bartlett_info.append(str_variables)
                                bartlett_info.append('')
                                str_chi=f'{chi:.4f}'
                                chi_item='Chi square value: ' + str_chi
                                print(chi_item)
                                bartlett_info.append(chi_item)
                                str_p=f'{p_val:.4f}'
                                p_item='p-value: ' + str_p
                                print(p_item)
                                bartlett_info.append(p_item)

                                print('NB: Valid if p-value is significant')
                                print('...')
                                bartlett_df=pd.DataFrame(bartlett_info)
                            
                                
                                if valFactor['factor_table']==True:
                                    engine = 'xlsxwriter'
                                    filename=valFactor['factor_filename']
                                    new_excel_file=os.path.join(out_folder, filename + '.xlsx')
                                    with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                        bartlett_df.to_excel(writer, sheet_name="Chi and p", index = False, header=False, startrow=0, startcol=0)
                                    print('Output saved to Excel file')
                                    print('')
                            
                            #KMO test
                            if valFactor['kmo']==True:
                                kmo_all,kmo_model=calculate_kmo(data)
                                #print('Hele')
                                #kmo_all
                                #kmo_model
                                str_kmo=f'{kmo_model}'
                                print('Value of KMO: ' + str_kmo)
                                print('NB: KMO < 0.6 is considered inadequate')
                                print('...')
                                
                                kmo_list=[]
                                kmo_list.append('KMO-values:')
                                kmo_list.append('')
                                kmo_list.append(str_variables)
                                kmo_list.append('')
                                kmo_list.append('Value of KMO: ' + str_kmo)
                                kmo_list.append('NB: KMO < 0.6 is considered inadequate')
                                kmo_df=pd.DataFrame(kmo_list)

                                if valFactor['factor_table']==True:
                                    engine = 'xlsxwriter'
                                    filename=valFactor['factor_filename']
                                    new_excel_file=os.path.join(out_folder, filename + '.xlsx')
                                    with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                        kmo_df.to_excel(writer, sheet_name="kmo", index = False, header=False, startrow=0, startcol=0)
                                    print('Output saved to Excel file')
                                    print('')

                            #Get number of factors
                            if valFactor['kai_crit']==True:
                                
                                print('All eigenvalues:')
                                eigenvalues_df=(pd.DataFrame(ev, index=fact_variables))
                                print('')
                                
                                print(eigenvalues_df)
                                
                                nr_eigenvalues_above_1=[]
                                for index, row in eigenvalues_df.iterrows():
                                    if float(row[0])>1:
                                        nr_eigenvalues_above_1.append(row)
                                
                                nr=len(nr_eigenvalues_above_1)
                                
                                print('') 
                                print('Number of factors should be number of eigenvalues > 1 (the Kaiser criterion).')
                                print('')
                                nr_str=f'{nr}'
                                print('Eigenvalues > 1 = ' + nr_str)
                                print('...')

                                kaicrit_list=[]
                                kaicrit_list.append('Eigenvalues:')
                                kaicrit_list.append('')
                                kaicrit_list.append(str_variables)
                                kaicrit_list.append('')
                                info_df=pd.DataFrame(kaicrit_list)

                                if valFactor['factor_table']==True:
                                    engine = 'xlsxwriter'
                                    
                                    with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                        info_df.to_excel(writer, sheet_name="Eigenvalues", index = False, header=False, startrow=0, startcol=0)
                                        eigenvalues_df.to_excel(writer, sheet_name="Eigenvalues", index = True, header=False, startrow=5, startcol=0)
                                    
                                    print('Eigenvalues saved to Excel file')
                                    print('')

                                if valFactor['factor_scree_png']==True: 
                                    plt.scatter(range(1,data.shape[1]+1),ev)
                                    plt.plot(range(1,data.shape[1]+1),ev)
                                    plt.title('Scree Plot')
                                    plt.xlabel('Factors')
                                    plt.ylabel('Eigenvalues')
                                    plt.grid()
                                    
                                    plt.savefig(os.path.join(out_folder, filename + '.png'), dpi=300, format='png', transparent=True)
                                    fig=None
                                    
                                    print('Scree plot saved as png file')
                                    print('...')

                                if valFactor['factor_scree_pdf']==True:  
                                    plt.scatter(range(1,data.shape[1]+1),ev)
                                    plt.plot(range(1,data.shape[1]+1),ev)
                                    plt.title('Scree Plot')
                                    plt.xlabel('Factors')
                                    plt.ylabel('Eigenvalue')
                                    plt.grid()
                                    
                                    plt.savefig(os.path.join(out_folder, filename + '.pdf'), dpi=300, format='pdf', transparent=True)
                                    fig=None
                                    
                                    print('Scree plot saved as PDF file')
                                    print('...')

                            #Confirmatory analysis
                            if valFactor['confirmatory']==True:
                                if valFactor['varimax']==True:
                                    rot="varimax"
                                elif valFactor['promax']==True:
                                    rot="promax"
                                elif valFactor['oblimax']==True:
                                    rot="oblimax"
                                elif valFactor['oblimin']==True:
                                    rot="oblimin"
                                else:
                                    rot=None

                                print(rot)

                                model_dict = factor_dict
                                model_spec = ModelSpecificationParser.parse_model_specification_from_dict(data, model_dict)
                                cfa = ConfirmatoryFactorAnalyzer(model_spec, disp=False) 
                                cfa.fit(data.values) 
                                #cfa.loadings_
                                cfa_loadings=cfa.loadings_
                                cfa_df=pd.DataFrame(cfa.loadings_, index=data.columns)
                                print('List:')
                                #column_names_list.insert(0, 'Variable_names:')
                                print(column_names_list)
                                cfa_df.columns = column_names_list

                                
                                print('Factor loadings:')
                                print('')
                                print(cfa_df)

                                #loadings_df=pd.DataFrame(loadings, index=fact_variables)
                                #print(loadings_df)
                                #print('...')

                                print('ok1')
                                if valFactor['factor_table']==True:
                                    list_info=[]
                                    list_info.append(f'Factor loadings, exploratory analysis, rotation = {rot}:')
                                    list_info.append('')
                                    list_info.append(str_variables)
                                    list_info.append('')
                                    df_info=pd.DataFrame(list_info)
                                    df_factornames = pd.DataFrame(column_names_list)
                                    if rot == None:
                                        rot = 'unrotated'
                                    
                                    new_excel_file=os.path.join(out_folder, filename + '_' + rot + '_CFA_max_li.xlsx')
                                    engine = 'xlsxwriter'
                                
                                    print('Saving...')
                                    with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                        df_info.to_excel(writer, sheet_name="Confirmatory factor loadings", index = False, header=False, startrow=0, startcol=0)
                                        #df_factornames.to_excel(writer, sheet_name="Confirmatory factor loadings", index = False, header=False, startrow=15, startcol=1)
                                        cfa_df.to_excel(writer, sheet_name="Confirmatory factor loadings", index = True, header=True, startrow=5, startcol=0)
                                    
                                    print('')
                                    print('CFA Factor loadings saved as Excel file.') 
                                    print('...')   
                                    

                            #Exploratory analysis
                            if valFactor['exploratory']==True:
        
                                n_fact=int(valFactor['nr_factors'])
                                n_fact_str=f"{valFactor['nr_factors']}"

                                list_factors=[]
                                n=1

                                while n < (n_fact + 1):
                                    factor=f'Factor {n}'
                                    list_factors.append(factor)
                                    n+=1 

                                #Varimax rotation
                                if valFactor['varimax']==True:
                                    rot="varimax"
                                    print('Rotation: ' + rot)
                                    print('Number of factors: ' + n_fact_str)
                                    fa=FactorAnalyzer(n_factors=n_fact, rotation=rot)
                                    fa.fit(data)
                                    rot_df=pd.DataFrame(fa.loadings_, index=data.columns)
                                    rot_df.columns=list_factors
                                    print('')
                                    print(rot_df)
                                    print('')
                                    print('Variance and communalities:')
                                    print('')
                                    rot_var=pd.DataFrame(fa.get_factor_variance(), index=['Variance','Proportional Var','Cumulative Var'])
                                    rot_var.columns=list_factors
                                    print(rot_var)
                                    print('')
                                    rot_com=(pd.DataFrame(fa.get_communalities(),index=data.columns,columns=['Communalities']))
                                    print(rot_com)
                                    print('')

                                    if valFactor['factor_table']==True:
                                        new_excel_file=os.path.join(out_folder, filename + '_varimax_minres.xlsx')
                                        engine = 'xlsxwriter'
                                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                            rot_df.to_excel(writer, sheet_name="Varimax rotation", index = True, header=True, startrow=0, startcol=0)
                                        print('Factor loadings saved to Excel file')
                                        print('')

                                #Promax rotation
                                elif valFactor['promax']==True:
                                    rot="promax"
                                    print('Rotation: ' + rot)
                                    print('Number of factors: ' + n_fact_str)
                                    fa=FactorAnalyzer(n_factors=n_fact, rotation=rot)
                                    fa.fit(data)
                                    rot_df=pd.DataFrame(fa.loadings_, index=data.columns)
                                    rot_df.columns=list_factors
                                    print('')
                                    print(rot_df)
                                    print('')
                                    print('Variance and communalities:')
                                    print('')
                                    rot_var=pd.DataFrame(fa.get_factor_variance(), index=['Variance','Proportional Var','Cumulative Var'])
                                    rot_var.columns=list_factors
                                    print(rot_var)
                                    print('')
                                    rot_com=(pd.DataFrame(fa.get_communalities(),index=data.columns,columns=['Communalities']))
                                    print(rot_com)
                                    print('')

                                    if valFactor['factor_table']==True:
                                        new_excel_file=os.path.join(out_folder, filename + '_promax_minres.xlsx')
                                        engine = 'xlsxwriter'
                                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                            rot_df.to_excel(writer, sheet_name="Promax rotation", index = True, header=True, startrow=0, startcol=0)
                                        print('Factor loadings saved to Excel file')
                                        print('')

                                #Oblimin rotation
                                elif valFactor['oblimin']==True:
                                    rot="oblimin"
                                    print('Rotation: ' + rot)
                                    print('Number of factors: ' + n_fact_str)
                                    fa=FactorAnalyzer(n_factors=n_fact, rotation=rot)
                                    fa.fit(data)
                                    rot_df=pd.DataFrame(fa.loadings_, index=data.columns)
                                    rot_df.columns=list_factors
                                    print('')
                                    print(rot_df)
                                    print('')
                                    print('Variance and communalities:')
                                    print('')
                                    rot_var=pd.DataFrame(fa.get_factor_variance(), index=['Variance','Proportional Var','Cumulative Var'])
                                    rot_var.columns=list_factors
                                    print(rot_var)
                                    print('')
                                    rot_com=(pd.DataFrame(fa.get_communalities(),index=data.columns,columns=['Communalities']))
                                    print(rot_com)
                                    print('')

                                    if valFactor['factor_table']==True:
                                        new_excel_file=os.path.join(out_folder, filename + '_oblimin_minres.xlsx')
                                        engine = 'xlsxwriter'
                                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                            rot_df.to_excel(writer, sheet_name="Oblimin rotation", index = True, header=True, startrow=0, startcol=0)
                                        print('Factor loadings saved to Excel file')
                                        print('')

                                #Oblimax rotation
                                elif valFactor['oblimax']==True:
                                    rot="oblimax"
                                    print('Rotation: ' + rot)
                                    print('Number of factors: ' + n_fact_str)
                                    fa=FactorAnalyzer(n_factors=n_fact, rotation=rot)
                                    fa.fit(data)
                                    rot_df=pd.DataFrame(fa.loadings_, index=data.columns)
                                    rot_df.columns=list_factors
                                    print('')
                                    print(rot_df)
                                    print('')
                                    print('Variance and communalities:')
                                    print('')
                                    rot_var=pd.DataFrame(fa.get_factor_variance(), index=[f"{'Variance':>20}",f"{'Proportional Var':>20}",f"{'Cumulative Var':>20}"])
                                    rot_var.columns=list_factors
                                    print(rot_var)
                                    print('')
                                    rot_com=(pd.DataFrame(fa.get_communalities(),index=data.columns,columns=['Communalities']))
                                    print(rot_com)
                                    print('')

                                    if valFactor['factor_table']==True:
                                        new_excel_file=os.path.join(out_folder, filename + '_oblimax_minres.xlsx')
                                        engine = 'xlsxwriter'
                                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                            rot_df.to_excel(writer, sheet_name="Oblimax rotation", index = True, header=True, startrow=0, startcol=0)
                                        print('Factor loadings saved to Excel file')
                                        print('')

                                #Without rotation
                                elif valFactor['none']==True:
                                    rot=None
                                    print('Rotation: None')
                                    print('Number of factors: ' + n_fact_str)
                                    fa=FactorAnalyzer(n_factors=n_fact, rotation=rot)
                                    fa.fit(data)
                                    rot_df=pd.DataFrame(fa.loadings_, index=data.columns)
                                    rot_df.columns=list_factors
                                    print('')
                                    print(rot_df)
                                    print('')
                                    print('Variance and communalities:')
                                    print('')
                                    rot_var=pd.DataFrame(fa.get_factor_variance(), index=[f"{'Variance':>}",f"{'Proportional Var':>}",f"{'Cumulative Var':>}"])
                                    rot_var.columns=list_factors
                                    print(rot_var)
                                    print('')
                                    rot_com=(pd.DataFrame(fa.get_communalities(),index=data.columns,columns=['Communalities']))
                                    print(rot_com)
                                    print('')

                                    if valFactor['factor_table']==True:
                                        new_excel_file=os.path.join(out_folder, filename + 'unrotated_minres.xlsx')
                                        engine = 'xlsxwriter'
                                        print('saving')
                                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                            rot_df.to_excel(writer, sheet_name="Without rotation", index = True, header=True, startrow=0, startcol=0)
                                            rot_var.to_excel(writer, sheet_name="Without rotation", index = True, header=True, startrow=(len(rot_df.index) + 2), startcol=0)
                                            rot_com.to_excel(writer, sheet_name="Without rotation", index = True, header=True, startrow=(len(rot_df.index) + len(rot_var.index) + 4), startcol=0)
                                        
                                        print('Factor loadings saved to Excel file')
                                        print('')

                        break                
                    #evFactor=''
        
        #Latent class analysis
        
        if (not winLca_active) and (valOriginal['lca']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winLca_active=True

            layoutLca = [
                [sg.Text('Latent Class Analysis:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('')],
                [sg.Text('Variables:', size=(13,1)), sg.InputText('', key='lca_variables', size=(52,3))],
                [sg.Text('Choose meassurement:', font=('bold'))],
                [sg.Radio("Continous: ", 'RADIO1', key='lca_continous', default=False, size=(13,1)), sg.Radio('Binary', 'RADIO1', key='lca_binary', default=False, size=(13,1)), sg.Radio('Categorical', 'RADIO1', key='lca_categorical', default=False, size=(13,1))], 
                [sg.Text('Number of latent classes: ', size=(20, 1)),sg.InputText('', key='nr_lca', size=(3, 1))],
                [sg.Text('', size=(13,1))],
                [sg.Text('Output (optional):', font=('bold'), size=(35, 1))],
                [sg.Checkbox('Export new dataset with classes', key='lca_table', default=False, size=(30,1))],
                [sg.Text('Output filename:', size=(13,1)), sg.InputText('', key='lca_filename', size=(32,1)), sg.Text('(NB: Excl. file extension)', size=(20,1))],
                [sg.Text('Output folder:', size=(13,1)), sg.InputText('', key='lca_output', size=(52,1)), sg.FolderBrowse()],
                [sg.Text('', size=(13,1)), sg.Text('NB: If left open, dataset folder will be used')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutLca, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('For all tests and analyses:', font=('bold'), size=(50,1), background_color='lightsteelblue2')],
                [sg.Text('Please enter either all included variables separated by commas, or enter excluded variables with a minus sign in front of the first variable. If left open, all variables in the dataset will be included.', size=(50,3), background_color='lightsteelblue2')]], size=(500,170), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(60,80))]], size=(505,430), background_color='lightsteelblue2')]]
            
            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteLca = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]


            winLca=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteLca)

            while True:
                evLca, valLca = winLca.Read(timeout=100)

                if evLca=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evLca==''
                
                if evLca is None or evLca=='Back' or evLca=='Step back and change dataset':
                    winLca_active=False
                    winLca.Close()
                    del winLca
                    winOriginal.UnHide()
                    break       

                #Valdidate variables supplied

                if(evLca=='Continue' and valLca['lca_continous']==False and valLca['lca_binary']==False and valLca['lca_categorical']==False):
                    popup_break(evLca, 'Please choose meassurment.')
                elif(evLca=='Continue' and (valLca['lca_table']==True) and valLca['lca_filename']==''):
                    popup_break(evLca, 'Please enter a name for the export file.')
                elif(evLca=='Continue' and valLca['nr_lca']==''):
                    popup_break(evLca, 'Please enter number of latent classes.')
                
                elif evLca=='Continue':
                    
                    while True:      
                        validvariables=list(data_df.columns)
                        
                        filename=valLca['lca_filename'] 
                        nr_lca = int(valLca['nr_lca'])

                        if valLca['lca_variables'] == '':
                            varsplit=validvariables
                            varsplit=validvariables.replace(', ', ',')
                            data=data_df
                            '''
                            #convert to numeric values
                            data = data.apply(pd.to_numeric, errors='coerce')
                            #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                            data = data.dropna(axis=0)
                            '''
                       
                        else:
                            var_=valLca['lca_variables']
                            if var_.startswith('-'):
                                var=var_[1:]
                                if var.endswith(','):
                                    var=var.rstrip(',')
                                if var.startswith(','):
                                    var=var.lstrip(',')
                                var=var.replace(', ', ',')
                                var=var.replace(' ','')
                                var=var.strip()
                                varsplit=var.split(',')
                                var=str(var)
                                varsplit=var.split(',')
                    
                                
                                data=data_df.drop(varsplit, axis=1)
                                '''
                                #convert to numeric values
                                #data = data.apply(pd.to_numeric, errors='coerce')
                                #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                                data = data.dropna(axis=0)
                                '''

                            else:
                                var=var_

                                if var.endswith(','):
                                    var=var.rstrip(',')
                                if var.startswith(','):
                                    var=var.lstrip(',')
                                var=var.replace(', ', ',')
                                var=var.replace(' ','')
                                var=var.strip()
                                
                                #varsplit=var.split(',')
                                var=str(var)
                                varsplit=var.split(",")
                                
                                data=data_df[varsplit]
     
                                '''
                                #convert to numeric values
                                #data = data.apply(pd.to_numeric, errors='coerce')
                                #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                                data = data.dropna(axis=0)
                                '''

                        #Convert to numeric
                        ##data=data.apply(pd.to_numeric, downcast='integer', errors='coerce')
                        #Use numeric values only
                        #data = data.apply(pd.to_numeric, errors='coerce')
                        #remove rows with missing values:
                        ''''
                        data = data.dropna(axis=0)
                        #Converts floats to ints
                        ##data = data.astype(int)
                        #pd.set_option('display.max_rows', None)
                        '''
                        lca_variables=list(data.columns)
                    
                        str_variables='Included varables: '
                        n=0
                        for str in lca_variables:
                            if n==0:
                                str_variables=str_variables + str
                                n=1
                            else:
                                str_variables=str_variables + ', ' + str
                        
                        nr_variables=len(lca_variables)

                        response='yes'

                        for string in varsplit:    
                            if string not in validvariables:
                                response='no'
                        
                        if response=='no':
                            popup_break(evLca, 'Error: One or more entered variables not in dataset')

                        #while True:
                        if response=='yes':
                    
                            if valLca['lca_output'] == '':
                                out_folder=os.path.dirname(data_in)
                            else:
                                out_folder=valLca['lca_output'] 
                            
                            if valLca['lca_continous'] == True:
                                lca_meassurement = "continous"
                            elif valLca['lca_binary'] == True:
                                lca_meassurement = "binary"  
                            elif valLca['lca_categorical'] == True:
                                lca_meassurement = "categorical"
                            
                            print(lca_meassurement)
                            # Gaussian mixture model
                            model = StepMix(n_components=nr_lca, measurement=lca_meassurement, verbose=1, random_state=123)
                            #Create text file
                            
                            # Fit to data
                            lca_fit = model.fit(data)

                            # Edit dataset and save:    
                            if len(data) == len(data_df):
                                lca_df_add = model.predict(data)
                                data_df['Latent_class_nr'] = lca_df_add
                            else:
                                data_df['Latent_class_nr'] = ''

                            if valLca['lca_table']==True:
                                
                                # Measurement parameters
                                mm = model.get_mm_df(data)
                               
                            
                                engine = 'xlsxwriter'
                                filename=valLca['lca_filename']
                                new_excel_file=os.path.join(out_folder, filename + '.xlsx')
                                with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                                    data_df.to_excel(writer, sheet_name="Edited dataset", index = False, header=True, startrow=0, startcol=0)
                                    mm.to_excel(writer, sheet_name="LCA Values", index = True, header=True, startrow=0, startcol=0)
                                    
                                print('Output saved to Excel file')
                                print('')
                            
                                

                        break                
                    
        

        #SAV - converter
        if (not winSav_active) and (valOriginal['sav']==True) and (evOriginal=='Continue'):
            winOriginal.Hide()
            winSav_active=True

            

            layoutSav= [
                [sg.Text('')], 
                [sg.Text('SAV-converter:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('')], 
                [sg.Text('Choose SAV-file:', font=('bold'))],
                [sg.In('', key='sav-file'), sg.FileBrowse()],
                [sg.Text('Select output folder:', size=(35, 1), font=('bold'))],      
                [sg.InputText('', key='sav_output'), sg.FolderBrowse()],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Button('Convert'), sg.Button('Back')]]

            
            
            layoutOriginalFrame = [
                [sg.Frame('Function', layoutSav, size= (570, 600))]]

            

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('', background_color='lightsteelblue2')],
                [sg.Text('"SAV-converter" converts sav-files to csv- and xlsx-files', background_color='lightsteelblue2', size = (50, 2), justification='left')], 
                [sg.Output(size=(60,8))]], background_color='lightsteelblue2', size=(500,600))]]
            
            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteSav = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winSav=sg.Window('CORals Analytics v. 3.12.0', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteSav)

            while True:
                evSav, valSav = winSav.Read(timeout=100)

                if evSav is None or evSav == 'Back':
                    winSav_active=False
                    winSav.Close()
                    del winSav
                    winOriginal.UnHide()
                    break
                
                elif (evSav=='Convert') and (valSav['sav-file']==''):
                    popup_break(evSav, 'Choose sav-file')
                
                elif (evSav=='Convert') and not (str(valSav['sav-file']).endswith('.sav')):
                    popup_break(evSav, 'File has to be a sav-file')

                elif (evSav=='Convert') and (valSav['sav_output']==''):
                    popup_break(evSav, 'Choose output folder')

                elif evSav=='Convert':   
                    while True:
                        
                        file=valSav['sav-file']               
                        filename=os.path.basename(file)
                        sav_file=str(filename)
                        print('')
                        print('Name of sourcefile:')
                        print(filename)
                        print('')
                        print('Names of converted files: ')
                        print(filename + '_converted.xlsx')
                        print(filename + '_converted.csv')
                        
                        folder=file.replace('/'+sav_file,'')

                        df, meta = py.read_sav(file)
                        
                        print('Converting ...')

                        engine = 'xlsxwriter'
                        new_excel_file=os.path.join(folder, filename) + '_converted' + '.xlsx'
                        with pd.ExcelWriter(new_excel_file, engine=engine) as writer:
                            df.to_excel(writer, sheet_name="Converted dataset", index = False, header=True, startrow=0, startcol=0)
                        
                        df.to_csv(os.path.join(folder, filename) + '_converted' + '.csv')

                        print('')
                        print('SAV-file successfully converted.')
                
                        break    
                           
        #Cronbach's Alpha
        
        if (not winCronbachs_active) and (valOriginal['Cronbachs']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winCronbachs_active=True


            layoutCronbachs = [
                [sg.Text('')],
                [sg.Text("Chronbach's Alpha:", size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('')],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('')],
                [sg.Text('Enter variables to be analyzed:', font=('bold'))],
                [sg.Text('(Variable names are the column headers (row 1) in the dataset file.)')],
                [sg.Text('Divide variables by commas only (no white space).')],
                [sg.InputText('', key='Cronbachs_variables')],
                [sg.Text('')],
                [sg.Text('Export output: ', font=('bold')), sg.Radio('No: ', 'RADIO1', key='Cronbachs_output_no', default=True, size=(20,1)), sg.Radio('Yes: ', 'RADIO1', key='Cronbachs_output_yes', default=False, size=(20,1))],
                [sg.Text('')],
                [sg.Text('Select output folder:', font=('bold'), size=(35, 1))],      
                [sg.InputText('', key='Cronbachs_output'), sg.FolderBrowse()],
                [sg.Text('')],
                [sg.Text('')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutCronbachs, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('"Cronbachs Alpha ....', size=(50,2), background_color='lightsteelblue2')], 
                [sg.Text('', size=(50,2), background_color='lightsteelblue2')]], size=(500,100), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(500,445))]], size=(505,500), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteCronbachs = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winCronbachs=sg.Window('CORals Analytics v. 3.12.0', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteCronbachs)

            while True:
                evCronbachs, valCronbachs = winCronbachs.Read(timeout=100)    

                if evCronbachs=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evCronbachs==''

                if evCronbachs is None or evCronbachs=='Back' or evCronbachs=='Step back and change dataset':
                    winCronbachs_active=False
                    winCronbachs.Close()
                    del winCronbachs
                    winOriginal.UnHide()
                    break       

                    
                elif (evCronbachs=='Continue' and valCronbachs['Cronbachs_variables']==''):
                    popup_break(evCronbachs, 'Choose variables')

                elif ((evCronbachs=='Continue' and valCronbachs['Cronbachs_output_yes']==True) and valCronbachs['Cronbachs_output'] == ''):
                    popup_break(evCronbachs, 'Choose output folder')


                #Valdidate variables supplied
                elif (evCronbachs=='Continue') and not (valCronbachs['Cronbachs_variables']==''):
                    var=valCronbachs['Cronbachs_variables']
                    var = ''.join(char for char in var if char in printable)
                    var=var.replace(" ", "")
                    varsplit=var.split(',')

                    data=data_df[varsplit]
                    #remove rows with empty cells
                    data = data.dropna(axis=0)
                    #Remove rows with non-integar data
                    #data = data[data.applymap(CheckInt).all(axis=1)].astype(int)
                    #Remove rows with non-numeric data
                    data = data.apply(pd.to_numeric, errors='coerce')

                    validvariables=list(data.columns)
                    
                    list_var=list(varsplit)
                    
                    var=valCronbachs['Cronbachs_variables']
                    
                    varsplit=var.split(',')

                    data=data_df[varsplit]
                    

                    list_var=list(varsplit)

                    now=datetime.today().strftime('%Y-%m-%d_%H-%M-%S')
    
                    response='yes'
                    for string in list_var:    
                        if string not in validvariables:
                            response='no'
                            
                    if (evCronbachs=='Continue') and (response =='no'):
                        popup_break(evCronbachs, 'Variable(s) not in dataset')

                    #End validation of variables
                    
              
                    elif (evCronbachs=='Continue' and valCronbachs['Cronbachs_output_yes']==True and valCronbachs['Cronbachs_output']==''):
                        popup_break(evCronbachs, 'Choose output folder')   

                    elif (evCronbachs=='Continue') and not (valCronbachs['Cronbachs_variables']=='') and (response=='yes'):
                        cronbachs_95=str(pg.cronbach_alpha(data=data))
                        cronbachs_99=str(pg.cronbach_alpha(data=data, ci=.99))
                        str_remove = ['(', '[', ')', ']', 'array']
                        for char in str_remove:
                            cronbachs_95=cronbachs_95.replace(char, '')
                            cronbachs_99=cronbachs_99.replace(char, '')
                        cronbachs_95_list=cronbachs_95.split(',')
                        cronbachs_99_list=cronbachs_99.split(',')

                        alpha = cronbachs_95_list[0]
                        min_95=cronbachs_95_list[1]
                        max_95=cronbachs_95_list[2]
                        min_99=cronbachs_99_list[1]
                        max_99=cronbachs_99_list[2]

                        cronbachs_df = pd.DataFrame(np.array([[alpha, min_95, max_95, min_99, max_99]]), columns=["Cronbach's Alpha", "CI 95% min", "CI 95% max", "CI 99% min", "CI 99% max"])

                        print(cronbachs_df)
                        print('')
                        print('Variables: ' + var)

                        var_df=pd.DataFrame(varsplit)

                        if (evCronbachs=='Continue' and valCronbachs['Cronbachs_output_yes']==True and valCronbachs['Cronbachs_output']!=''):
                            output_folder=valCronbachs['Cronbachs_output']

                                    
                            engine = 'xlsxwriter'
                            with pd.ExcelWriter(os.path.join(output_folder, 'cronbachs_alpha_') + now + '.xlsx', engine=engine) as writer:
                                cronbachs_df.to_excel(writer, sheet_name="Chronbach's alpha", index = None, header=True, startrow=1)
                                var_df.to_excel(writer, sheet_name="Chronbach's alpha", index = None, header=False, startrow=4)
                            
                            print('')
                            print("Cronbach's Alpha saved to file")

                    
                                   
                    
        #Documentation

        
        if (not winDocumentation_active) and (evOriginal=='Documentation'):
            winOriginal.Hide()
            winDocumentation_active=True

            
            layoutDoc1 = [
                [sg.Text("Documentation:", justification="left", font=("Arial", 20), background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],                
                [sg.Text("Framework:", font=("bold", 18), background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("Python version: 3.9. ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("All tools: Pandas 1.4.1.", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Button("Back")],
                [sg.Text("", background_color='lightsteelblue2')]]

            layoutDoc2 = [
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("Python tools/packages per task:", font=("bold", 18), background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')],
                [sg.Text("Distribution analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Matplotlib 3.5.1 and Seaborn 0.11.2.", background_color='lightsteelblue2')],                
                [sg.Text("Correlation analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Matplotlib 3.5.1 and Seaborn 0.11.2.", background_color='lightsteelblue2')],
                [sg.Text("Regression analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Matplotlib 3.5.1 and Statsmodels 0.13.2.", background_color='lightsteelblue2')],
                [sg.Text("Mediation analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Pingouin 0.5.1.", background_color='lightsteelblue2')],
                [sg.Text("Moderation analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("PyProcessMacro 1.0.12 and Matplotlib 3.5.1.", background_color='lightsteelblue2')],
                [sg.Text("Factor analysis: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Factor-analyzer 0.4.0 and Matplotlib 3.5.1.", background_color='lightsteelblue2')],
                [sg.Text("CSV-rescue: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Cronbachs Alpha: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Pingouin 0.5.1.", background_color='lightsteelblue2')],
                [sg.Text("SAV converter: ", font=("bold"), background_color='lightsteelblue2')],
                [sg.Text("Pyreadstat 1.1.6'.", background_color='lightsteelblue2')],
                [sg.Text("", background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompleteDoc = [
                [sg.Column(layoutDoc1, element_justification='left', size=(525,600), background_color='lightsteelblue2'), sg.Column(layoutDoc2, element_justification='left', size=(525,600), background_color='lightsteelblue2')]]

            layoutCompleteDocumentation = [
                [sg.Frame('Framework and tools/packages:', layoutCompleteDoc, size=(1050, 600), background_color='lightsteelblue2')], 
                [layoutBottom]]


            winDocumentation=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompleteDocumentation)

            while True:
                evDocumentation, valDocumentation = winDocumentation.Read(timeout=100)               

                if evDocumentation is None or evDocumentation=='Back':
                    winDocumentation_active=False
                    winDocumentation.Close()
                    del winDocumentation
                    winOriginal.UnHide()
                    break

        #Presentation

        if (not winPresentation_active) and (valOriginal['Presentation']==True) and (evOriginal=='Continue') and ((data_in.endswith('.csv')) or (data_in.endswith('.xlsx'))):
            winOriginal.Hide()
            winPresentation_active=True
            
            layoutPresentation = [
                [sg.Text('Make presentation:', size=(25,1), justification='left', font=("Arial", 20))],
                [sg.Text('Selected dataset:', font=('bold'))],
                [sg.Text(chosen_filename)],
                [sg.Button('Check available variables'), sg.Button('Step back and change dataset')],
                [sg.Text('', font=('bold'))],
                [sg.Text('Choose category:', size=(25,1)), sg.InputText('', key='new_variable', size=(15, 1))],
                [sg.Radio('Single variables: ', 'RADIO1', key='singles', default=False, size=(10,1)), sg.Radio('Index variables:', 'RADIO1', key='index', default=False, size=(20,1))],
                [sg.Text('Enter included variables, divided by comma')],
                [sg.InputText('', key='presentation_variables', size=(35,1))],  
                [sg.Text('Select output folder:', font=('bold'), size=(60, 1))],      
                [sg.InputText('', key='presentation_output'), sg.FolderBrowse()],
                [sg.Text('NB: If left open, source folder will be used')],
                [sg.Text('Enter name of presentation file (excl. ".pptx"):', font=('bold'), size=(35, 1))],      
                [sg.InputText('', key='presentation_filename', size=(35, 1)), sg.Text('NB: If left open, original filename will be used')],
                [sg.Button('Continue'), sg.Button('Back')]]

            layoutOriginalFrame = [
                [sg.Frame('Function', layoutPresentation, size= (570, 600))]]

            layoutRight = [
                [sg.Frame('Description', [
                [sg.Text('This tool creates a Powerpoint presentation with mean results', font=('bold'), size = (50, 1), justification='left', background_color='lightsteelblue2')], 
                [sg.Text('"Single variables" creates one slide for each variable.', size=(50,1), background_color='lightsteelblue2')],
                [sg.Text('"Index variable" creates one slide with results both for the index and for the included variables.', size=(55,1), background_color='lightsteelblue2')]], size=(500,150), background_color='lightsteelblue2')],
                [sg.Frame('Preview/log', [
                [sg.Output(background_color='white', size=(60,50))]], size=(505,450), background_color='lightsteelblue2')]]

            layoutBottom =  [
                [sg.Text('')],
                [sg.Text('©Christian Otto Ruge, licenced under GNU GPL v3')]]

            layoutCompletePresentation = [[sg.Column(layoutOriginalFrame, element_justification='left'), sg.Column(layoutRight, element_justification='right')], [layoutBottom]]

            winPresentation=sg.Window('CORals Analytics v. 3.12.2', default_element_size=(40, 1), grab_anywhere=False, location=(100,100), size=(1080,690)).Layout(layoutCompletePresentation)

            while True:
                evPresentation, valPresentation = winPresentation.Read(timeout=100)    

                if evPresentation=='Check available variables':
                    variables=list(data_df.columns)
                    separator = ','
                    print('Available variables in dataset:')
                    print(separator.join(variables))
                    print('')
                    evPresentation==''
                    

                if evPresentation is None or evPresentation=='Back' or evPresentation=='Step back and change dataset':
                    winPresentation_active=False
                    winPresentation.Close()
                    del winPresentation
                    winOriginal.UnHide()
                    break       

                    
                elif (evPresentation=='Continue' and (valPresentation['presentation_variables']=='')):
                    popup_break(evPresentation, 'Choose variables')


                #Valdidate variables supplied
                elif (evPresentation=='Continue'):

                    now = datetime.now()
                    now=now.strftime('%d/%m/%Y %H:%M:%S')

                    validvariables=list(data_df.columns)

                    while True:
                        response='yes'
                        var=valPresentation['presentation_variables']
                        var=var.replace(" ", "")
                        var=str(var)
                        varsplit=var.split(",")

                        for string in varsplit:    
                                if string not in validvariables:
                                    response='no' 
                            
                        if (evPresentation=='Continue') and (response =='no'):
                            popup_break(evPresentation, 'Selected variable(s) not in dataset')    

                        if valPresentation['singles']==True:
                            print('Creating single slides')
                            
                        elif valPresentation['index']==True:
                            print('Creating index slide')  

                        if valPresentation['presentation_filename'] == '':
                            new_filename=chosen_filename.replace('.xlsx', '.pptx')
                        else:
                            new_filename=valPresentation['presentation_filename'] + '.pptx'

                        if valPresentation['presentation_output'] == '':
                            out_folder=os.path.dirname(data_in)
                        else:
                            out_folder=valPresentation['presentation_output'] 

                        
                        
                        
                        
                        print('')
                        print('Presentation created, dataset saved')

                        

                        break
      
                            
except:
    sg.Popup('Ooops! Something went wrong! This may be due to invalid input in form fields. It may also be caused by the dataset file, or a file with the same name as the output file, may be open in another program (i.e. Excel). PLEASE CLOSE THIS WINDOW, close open files and retry running the program. If the problem persists, please feel free to contact CORals for support at www.corals.no/kontakt.')
    